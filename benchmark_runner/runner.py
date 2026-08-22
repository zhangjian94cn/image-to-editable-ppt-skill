"""Run frozen slide images through one local Codex task per page.

This is development tooling, not a product controller.  It deliberately owns
only corpus selection, process evidence, PowerPoint rendering, compact visual
diagnostics, and the snapshot-v2 projection requested by reviewers.
"""

from __future__ import annotations

import argparse
import hashlib
import json
import os
import re
import shutil
import subprocess
import sys
import time
from dataclasses import asdict, dataclass, field
from datetime import datetime, timezone
from pathlib import Path
from typing import Any, Iterable

import numpy as np
from PIL import Image, ImageChops, ImageFilter
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

# The benchmark is intentionally runnable from a clean checkout without first
# installing the Skill package into the caller's Python environment.
CHECKOUT_ROOT = Path(__file__).resolve().parents[1]
CLI_ROOT = CHECKOUT_ROOT / "skills/image-to-editable-ppt/cli"
if str(CLI_ROOT) not in sys.path:
    sys.path.insert(0, str(CLI_ROOT))

from editppt.source_space import prepare_authoring_source


ROOT_FILES = {"source.png", "candidate.pptx", "candidate.png", "report.md", "artifacts"}
DIRECT_POWERPOINT_PATTERN = re.compile(
    r"(?:osascript|cua-driver|pgrep|pkill|killall|\bopen\b).{0,240}(?:powerpoint|microsoft powerpoint)|"
    r"(?:powerpoint|microsoft powerpoint).{0,240}(?:osascript|cua-driver|pgrep|pkill|killall|\bopen\b)",
    re.IGNORECASE | re.DOTALL,
)
PARENT_SCAN_PATTERN = re.compile(r"\b(?:find|rg|grep)\b[^\n]{0,240}(?:\.\./|\s\.\.(?:\s|$))")
EXTERNAL_CONTEXT_PATTERN = re.compile(r"(?:/|\\)\.codex/(?:memories|memory)|\bMEMORY\.md\b", re.IGNORECASE)


class BenchmarkError(RuntimeError):
    """The benchmark cannot produce trustworthy evidence."""


@dataclass(frozen=True)
class PageCase:
    deck_id: str
    slide_number: int
    source: Path
    expected: Path

    @property
    def page_id(self) -> str:
        return f"{self.deck_id}-p{self.slide_number:03d}"


@dataclass
class CommandEvidence:
    command: str
    exit_code: int | None
    elapsed_sec: float | None = None


@dataclass
class PageOutcome:
    page_id: str
    verdict: str = "execution_failed"
    elapsed_sec: float = 0.0
    error: str = ""
    issues: list[dict[str, str]] = field(default_factory=list)
    metrics: dict[str, Any] = field(default_factory=dict)
    commands: list[CommandEvidence] = field(default_factory=list)
    skill_commands: list[CommandEvidence] = field(default_factory=list)
    decisions: list[str] = field(default_factory=list)


def _read_json(path: Path) -> dict[str, Any]:
    try:
        value = json.loads(path.read_text(encoding="utf-8"))
    except (OSError, json.JSONDecodeError):
        return {}
    return value if isinstance(value, dict) else {}


def _write_json(path: Path, value: Any) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    temporary = path.with_suffix(path.suffix + ".tmp")
    temporary.write_text(
        json.dumps(value, ensure_ascii=False, indent=2, sort_keys=True) + "\n",
        encoding="utf-8",
    )
    temporary.replace(path)


def _sha256(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as handle:
        for block in iter(lambda: handle.read(1024 * 1024), b""):
            digest.update(block)
    return digest.hexdigest()


def _git_value(repo: Path, *args: str) -> str:
    completed = subprocess.run(
        ["git", "-C", str(repo), *args], capture_output=True, text=True, check=False
    )
    return completed.stdout.strip() if completed.returncode == 0 else ""


def _now() -> str:
    return datetime.now(timezone.utc).replace(microsecond=0).isoformat()


def _run_id(label: str) -> str:
    safe = re.sub(r"[^0-9A-Za-z._-]+", "-", label.strip()).strip("-.") or "run"
    return datetime.now(timezone.utc).strftime("%Y%m%dT%H%M%SZ-") + safe


def _slug(value: str) -> str:
    normalized = re.sub(r"[^0-9A-Za-z._\-\u4e00-\u9fff]+", "-", value.strip())
    return re.sub(r"-+", "-", normalized).strip("-.") or "image-case"


def ingest_image(args: argparse.Namespace) -> int:
    corpus = Path(args.corpus).expanduser().resolve()
    source = Path(args.source).expanduser().resolve()
    hints = Path(args.text_hints).expanduser().resolve() if args.text_hints else None
    if not source.is_file():
        raise BenchmarkError(f"source image is unavailable: {source}")
    digest = _sha256(source)
    deck_id = f"{_slug(args.name)}-{digest[:12]}"
    deck_dir = corpus / "decks" / deck_id
    if deck_dir.exists():
        print(deck_id)
        return 0
    staging = corpus / ".staging" / f"{deck_id}-{os.getpid()}"
    page_dir = staging / "slides/001"
    page_dir.mkdir(parents=True, exist_ok=False)
    shutil.copy2(source, page_dir / "input.png")
    hint_payload = _read_json(hints) if hints else {}
    lines = hint_payload.get("lines") if isinstance(hint_payload.get("lines"), list) else []
    expected_objects: list[dict[str, Any]] = []
    for index, item in enumerate(lines, start=1):
        if not isinstance(item, dict) or not str(item.get("text") or "").strip():
            continue
        expected_objects.append({
            "shape_id": str(item.get("id") or f"P{index:02d}"),
            "name": str(item.get("id") or f"P{index:02d}"),
            "kind": "text", "text": str(item["text"]),
            "box_px": list(item.get("box_px") or []), "visible": True,
            "recognition_source": str(hint_payload.get("backend") or hint_payload.get("source") or "frozen-hints"),
        })
    with Image.open(source) as image:
        width, height = image.size
    _write_json(page_dir / "expected.json", {
        "schema_version": 1, "deck_id": deck_id, "slide_number": 1,
        "truth_mode": "frozen_image", "source": "input.png",
        "image_size": [width, height], "objects": expected_objects,
        "counts": {"text": len(expected_objects), "total": len(expected_objects)},
        "ocr": {"backend": str(hint_payload.get("backend") or "frozen-hints"),
                "region_count": len(expected_objects), "source_sha256": digest},
    })
    _write_json(staging / "deck.json", {
        "schema_version": 1, "deck_id": deck_id, "truth_mode": "frozen_image",
        "note": args.note, "original_sha256": digest,
        "slides": [{"number": 1, "path": "slides/001", "hidden": False}],
    })
    _write_json(staging / "evidence/ingest.json", {
        "schema_version": 1, "source_sha256": digest, "source_name": source.name,
        "text_hints_sha256": _sha256(hints) if hints and hints.is_file() else "",
        "ingested_at": _now(),
    })
    deck_dir.parent.mkdir(parents=True, exist_ok=True)
    staging.replace(deck_dir)
    index_path = corpus / "index.json"
    index = _read_json(index_path)
    decks = index.get("decks") if isinstance(index.get("decks"), list) else []
    decks.append({
        "deck_id": deck_id, "original_sha256": digest,
        "path": f"decks/{deck_id}", "reference_sha256": None,
        "slide_count": 1, "truth_mode": "frozen_image",
    })
    index["schema_version"] = int(index.get("schema_version") or 1)
    index["decks"] = sorted(decks, key=lambda value: str(value.get("deck_id") or ""))
    _write_json(index_path, index)
    print(deck_id)
    return 0


def _resolve_cases(corpus: Path, suite_name: str) -> list[PageCase]:
    suites = _read_json(corpus / "suites.json")
    raw: list[Any] | None = None
    page_suites = suites.get("page_suites")
    if isinstance(page_suites, dict):
        candidate = page_suites.get(suite_name)
        if isinstance(candidate, list):
            raw = candidate
    if raw is None:
        legacy = suites.get("suites")
        candidate = legacy.get(suite_name) if isinstance(legacy, dict) else None
        if not isinstance(candidate, list):
            raise BenchmarkError(f"unknown benchmark suite: {suite_name}")
        raw = candidate

    selected: list[tuple[str, int | None]] = []
    for item in raw:
        if isinstance(item, str):
            selected.append((item, None))
        elif isinstance(item, dict) and item.get("deck_id"):
            slide = item.get("slide")
            selected.append((str(item["deck_id"]), int(slide) if slide is not None else None))
        else:
            raise BenchmarkError(f"invalid suite entry: {item!r}")

    cases: list[PageCase] = []
    seen: set[tuple[str, int]] = set()
    for deck_id, requested_slide in selected:
        deck_dir = corpus / "decks" / deck_id
        deck = _read_json(deck_dir / "deck.json")
        if not deck:
            raise BenchmarkError(f"missing deck metadata: {deck_id}")
        slides = deck.get("slides") if isinstance(deck.get("slides"), list) else []
        for slide in slides:
            if not isinstance(slide, dict):
                continue
            number = int(slide.get("number") or 0)
            if requested_slide is not None and number != requested_slide:
                continue
            key = (deck_id, number)
            if key in seen:
                continue
            page_dir = deck_dir / "slides" / f"{number:03d}"
            source = page_dir / "input.png"
            expected = page_dir / "expected.json"
            if not source.is_file():
                raise BenchmarkError(f"missing frozen source: {source}")
            cases.append(PageCase(deck_id, number, source, expected))
            seen.add(key)
        if requested_slide is not None and (deck_id, requested_slide) not in seen:
            raise BenchmarkError(f"suite requests missing page: {deck_id}/{requested_slide}")
    if not cases:
        raise BenchmarkError(f"suite is empty: {suite_name}")
    return cases


def _filter_cases(cases: list[PageCase], requested: list[str] | None) -> list[PageCase]:
    if not requested:
        return cases
    wanted = set(requested)
    selected = [case for case in cases if case.page_id in wanted]
    missing = wanted - {case.page_id for case in selected}
    if missing:
        raise BenchmarkError(f"requested pages are not in the suite: {sorted(missing)}")
    return selected


def _page_prompt(skill_root: Path) -> str:
    template = skill_root / "prompts/page-task.md"
    if template.is_file():
        text = template.read_text(encoding="utf-8")
        return text.replace("{{SKILL_ROOT}}", str(skill_root))
    return f"""Use $image-to-editable-ppt to rebuild source.png as one object-level editable PowerPoint slide.

Read and follow the installed Skill at {skill_root / 'SKILL.md'}.
Work only in the current directory. Do not create subagents or controller state.
Do not place source.png or an almost-whole-page crop into the slide.
Generate page.pptx and then write result.json with status=ready and output_pptx=page.pptx.
"""


def _codex_command(args: argparse.Namespace, work: Path, source: Path) -> list[str]:
    codex = shutil.which(args.codex_bin) if not Path(args.codex_bin).is_file() else args.codex_bin
    if not codex:
        raise BenchmarkError(f"Codex CLI is unavailable: {args.codex_bin}")
    command = [
        str(codex), "exec", "--profile", args.profile, "--json",
        "--skip-git-repo-check", "-C", str(work), "-s", "danger-full-access",
    ]
    if args.model:
        command += ["-m", args.model]
    if args.effort:
        command += ["-c", f'model_reasoning_effort="{args.effort}"']
    command += [
        "--image", str(source), "--output-last-message",
        str(work.parent.parent / "codex" / "last-message.txt"), "-",
    ]
    return command


def _execute_codex(
    args: argparse.Namespace, skill_root: Path, work: Path, codex_dir: Path,
) -> tuple[int, float, str]:
    source = work / "source.png"
    command = _codex_command(args, work, source)
    codex_dir.mkdir(parents=True, exist_ok=True)
    env = os.environ.copy()
    env["MIAOBI_SKILL_ROOT"] = str(skill_root)
    env["EDITPPT_TRACE_FILE"] = str(work.parent.parent / "telemetry.jsonl")
    skill_bin = skill_root / "cli/.venv/bin"
    if skill_bin.is_dir():
        env["PATH"] = str(skill_bin) + os.pathsep + env.get("PATH", "")
    started = time.monotonic()
    with (codex_dir / "events.jsonl").open("w", encoding="utf-8") as events, (
        codex_dir / "stderr.log"
    ).open("w", encoding="utf-8") as errors:
        process = subprocess.Popen(
            command, stdin=subprocess.PIPE, stdout=events, stderr=errors,
            text=True, env=env, start_new_session=True,
        )
        assert process.stdin is not None
        process.stdin.write(_page_prompt(skill_root))
        process.stdin.close()
        returncode = process.wait()
    elapsed = time.monotonic() - started
    stderr = (codex_dir / "stderr.log").read_text(encoding="utf-8", errors="replace")
    return returncode, elapsed, stderr[-4000:]


def _render_powerpoint(
    candidate: Path, render_dir: Path, skill_root: Path,
) -> tuple[Path | None, dict[str, Any], float, str]:
    render_dir.mkdir(parents=True, exist_ok=True)
    editppt = skill_root / "cli/.venv/bin/editppt"
    if not editppt.is_file():
        resolved = shutil.which("editppt")
        if not resolved:
            return None, {}, 0.0, "Skill editppt executable is unavailable"
        editppt = Path(resolved)
    rendered = render_dir / "candidate.png"
    command = [
        str(editppt), "render", str(candidate.parent), "--input", candidate.name,
        "--out", str(rendered), "--evidence-dir", str(render_dir), "--dpi", "200",
    ]
    started = time.monotonic()
    completed = subprocess.run(command, capture_output=True, text=True, check=False)
    elapsed = time.monotonic() - started
    report = _read_json(render_dir / "powerpoint-render.json")
    if not report:
        try:
            value = json.loads(completed.stdout)
            report = value if isinstance(value, dict) else {}
        except json.JSONDecodeError:
            report = {}
    if completed.returncode != 0 or report.get("success") is not True or not rendered.is_file():
        error = str(report.get("error") or completed.stderr.strip() or "PowerPoint render failed")
        return None, report, elapsed, error
    return rendered, report, elapsed, ""


def _normal_text(value: str) -> str:
    normalized = str(value or "").translate(
        str.maketrans({"丨": "|", "｜": "|", "︱": "|"})
    )
    return re.sub(r"[\s\u000b]+", "", normalized).casefold()


def _box_overlap(first: Any, second: Any) -> float:
    """Return how much of ``first`` is covered by ``second`` in source pixels."""

    if not isinstance(first, list) or not isinstance(second, list) or len(first) != 4 or len(second) != 4:
        return 0.0
    left, top, width, height = [float(value) for value in first]
    other_left, other_top, other_width, other_height = [float(value) for value in second]
    if width <= 0 or height <= 0 or other_width <= 0 or other_height <= 0:
        return 0.0
    overlap_width = max(0.0, min(left + width, other_left + other_width) - max(left, other_left))
    overlap_height = max(0.0, min(top + height, other_top + other_height) - max(top, other_top))
    return overlap_width * overlap_height / (width * height)


def _expected_text_evidence(expected: dict[str, Any]) -> tuple[list[str], list[dict[str, str]]]:
    """Select rendered text, excluding placeholders and source objects hidden by screenshots."""

    objects = expected.get("objects") or []
    values: list[str] = []
    excluded: list[dict[str, str]] = []
    for index, item in enumerate(objects):
        if not isinstance(item, dict) or item.get("kind") != "text" or item.get("visible", True) is False:
            continue
        value = str(item.get("text") or "")
        compact = _normal_text(value)
        if re.fullmatch(r"[‹<«\[]?#(?:›|>|»|\])?", compact):
            excluded.append({"text": value, "reason": "dynamic_placeholder"})
            continue
        hidden_by_picture = any(
            isinstance(later, dict)
            and later.get("kind") == "picture"
            and _box_overlap(item.get("box_px"), later.get("box_px")) >= 0.95
            for later in objects[index + 1 :]
        )
        if hidden_by_picture:
            excluded.append({"text": value, "reason": "occluded_by_later_picture"})
            continue
        if compact:
            values.append(value)
    return values, excluded


def _pptx_metrics(candidate: Path, expected_path: Path) -> dict[str, Any]:
    presentation = Presentation(candidate)
    if len(presentation.slides) != 1:
        raise BenchmarkError(f"page candidate has {len(presentation.slides)} slides")
    slide = presentation.slides[0]
    width = float(presentation.slide_width)
    height = float(presentation.slide_height)
    text_values: list[str] = []
    pictures = tables = text_shapes = native_shapes = 0
    max_picture_coverage = 0.0
    out_of_bounds = 0
    for shape in slide.shapes:
        left, top = float(shape.left), float(shape.top)
        shape_width, shape_height = float(shape.width), float(shape.height)
        if left < -1 or top < -1 or left + shape_width > width + 1 or top + shape_height > height + 1:
            out_of_bounds += 1
        if getattr(shape, "has_text_frame", False):
            value = str(shape.text or "").strip()
            if value:
                text_shapes += 1
                text_values.append(value)
        if getattr(shape, "has_table", False):
            tables += 1
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            pictures += 1
            max_picture_coverage = max(
                max_picture_coverage,
                (shape_width * shape_height) / max(1.0, width * height),
            )
        elif shape.shape_type in {
            MSO_SHAPE_TYPE.AUTO_SHAPE, MSO_SHAPE_TYPE.FREEFORM,
            MSO_SHAPE_TYPE.GROUP, MSO_SHAPE_TYPE.LINE,
        }:
            native_shapes += 1

    expected = _read_json(expected_path)
    expected_text, excluded_expected_text = _expected_text_evidence(expected)
    candidate_text = _normal_text("\n".join(text_values))
    missing_text = [value for value in expected_text if _normal_text(value) not in candidate_text]
    matched = len(expected_text) - len(missing_text)
    coverage = matched / len(expected_text) if expected_text else None
    return {
        "slide_count": len(presentation.slides),
        "object_count": len(slide.shapes),
        "text_shape_count": text_shapes,
        "native_shape_count": native_shapes,
        "table_count": tables,
        "picture_count": pictures,
        "max_picture_coverage": round(max_picture_coverage, 6),
        "out_of_bounds_count": out_of_bounds,
        "expected_text_count": len(expected_text),
        "matched_text_count": matched,
        "missing_text_count": len(missing_text),
        "missing_texts": missing_text,
        "excluded_expected_text_count": len(excluded_expected_text),
        "excluded_expected_texts": excluded_expected_text,
        "text_coverage": round(coverage, 6) if coverage is not None else None,
    }


def _visual_metrics(source_path: Path, candidate_path: Path, compare_dir: Path) -> dict[str, Any]:
    compare_dir.mkdir(parents=True, exist_ok=True)
    with Image.open(source_path) as source_image, Image.open(candidate_path) as candidate_image:
        source = source_image.convert("RGB")
        candidate = candidate_image.convert("RGB").resize(source.size, Image.Resampling.LANCZOS)
    small_source = source.resize((128, max(1, round(source.height * 128 / source.width))), Image.Resampling.LANCZOS)
    small_candidate = candidate.resize(small_source.size, Image.Resampling.LANCZOS)
    first = np.asarray(small_source, dtype=np.float32)
    second = np.asarray(small_candidate, dtype=np.float32)
    coarse = float(np.abs(first - second).mean() / 255.0)
    source_gray = np.asarray(source.convert("L"), dtype=np.int16)
    candidate_gray = np.asarray(candidate.convert("L"), dtype=np.int16)
    # Treat only meaningful dark/color content as ink.  Very light swimlane
    # bands and shadows are visual styling, not missing semantic content.
    source_ink = source_gray < 210
    candidate_ink = candidate_gray < 210
    # A two-pixel tolerance removes anti-aliasing and tiny font-rendering
    # shifts while still exposing missing labels, lines, and regions.
    source_near = np.asarray(
        Image.fromarray((source_ink * 255).astype(np.uint8)).filter(ImageFilter.MaxFilter(5))
    ) > 0
    candidate_near = np.asarray(
        Image.fromarray((candidate_ink * 255).astype(np.uint8)).filter(ImageFilter.MaxFilter(5))
    ) > 0
    ink_loss = float(
        np.logical_or(source_ink & ~candidate_near, candidate_ink & ~source_near).mean()
    )
    difference = ImageChops.difference(source, candidate)
    difference.save(compare_dir / "diff.png")
    heat = np.asarray(difference.convert("L"), dtype=np.uint8)
    heatmap = np.zeros((source.height, source.width, 3), dtype=np.uint8)
    heatmap[..., 0] = np.clip(heat * 3, 0, 255)
    heatmap[..., 1] = np.clip(heat // 3, 0, 255)
    Image.fromarray(heatmap, mode="RGB").save(compare_dir / "heatmap.png")
    return {"coarse_rgb_loss": round(coarse, 6), "content_ink_loss": round(ink_loss, 6)}


def _events(path: Path) -> tuple[list[CommandEvidence], list[str], list[str]]:
    commands: list[CommandEvidence] = []
    decisions: list[str] = []
    changed: list[str] = []
    if not path.is_file():
        return commands, decisions, changed
    for line in path.read_text(encoding="utf-8", errors="replace").splitlines():
        try:
            event = json.loads(line)
        except json.JSONDecodeError:
            continue
        if event.get("type") != "item.completed":
            continue
        item = event.get("item") if isinstance(event.get("item"), dict) else {}
        item_type = item.get("type")
        if item_type == "command_execution":
            commands.append(CommandEvidence(str(item.get("command") or ""), item.get("exit_code")))
        elif item_type == "agent_message":
            text = str(item.get("text") or "").strip()
            if text and text not in decisions:
                decisions.append(text)
        elif item_type == "file_change":
            for value in item.get("changes") or []:
                if isinstance(value, dict) and value.get("path"):
                    changed.append(str(value["path"]))
    return commands, decisions, changed


def _skill_trace(*paths: Path) -> list[CommandEvidence]:
    evidence: list[CommandEvidence] = []
    seen: set[tuple[str, int | None, float | None]] = set()
    for path in paths:
        if not path.is_file():
            continue
        for line in path.read_text(encoding="utf-8", errors="replace").splitlines():
            try:
                value = json.loads(line)
            except json.JSONDecodeError:
                continue
            if not isinstance(value, dict):
                continue
            command = " ".join(
                str(item) for item in (value.get("command"), value.get("subcommand")) if item
            )
            item = CommandEvidence(
                command=f"editppt {command}".strip(),
                exit_code=value.get("exit_code"),
                elapsed_sec=float(value["elapsed_sec"]) if value.get("elapsed_sec") is not None else None,
            )
            key = (item.command, item.exit_code, item.elapsed_sec)
            if key not in seen:
                evidence.append(item)
                seen.add(key)
    return evidence


def _issues(
    metrics: dict[str, Any],
    render_error: str = "",
    commands: Iterable[CommandEvidence] = (),
    *,
    page_dir: Path | None = None,
    pages_root: Path | None = None,
) -> tuple[str, list[dict[str, str]]]:
    issues: list[dict[str, str]] = []
    if render_error:
        issues.append({"severity": "P0", "category": "execution", "message": render_error})
        return "execution_failed", issues
    if metrics.get("codex_powerpoint_rendered") is not True:
        issues.append({
            "severity": "P0",
            "category": "execution",
            "message": "Codex 未留下与最终 page.pptx SHA 绑定的 editppt render 证据",
        })
    forbidden = [value.command for value in commands if DIRECT_POWERPOINT_PATTERN.search(value.command)]
    if forbidden:
        issues.append({
            "severity": "P0",
            "category": "execution",
            "message": "页面 Codex 绕过 editppt render 直接控制或检查了 PowerPoint",
        })
    if page_dir is not None and pages_root is not None:
        current = str(page_dir.resolve())
        root = str(pages_root.resolve())
        cross_page = [
            value.command for value in commands
            if PARENT_SCAN_PATTERN.search(value.command)
            or (root in value.command and current not in value.command)
        ]
        if cross_page:
            issues.append({
                "severity": "P1",
                "category": "execution",
                "message": "页面 Codex 读取或扫描了父目录/其他页面证据，破坏单页独立性",
            })
    if any(EXTERNAL_CONTEXT_PATTERN.search(value.command) for value in commands):
        issues.append({
            "severity": "P1",
            "category": "execution",
            "message": "页面 Codex 读取了历史 memory，而不是仅依据当前源图和固定 Skill",
        })
    coverage = metrics.get("text_coverage")
    missing_texts = metrics.get("missing_texts") or []
    missing_summary = ""
    if missing_texts:
        excerpts = [re.sub(r"\s+", " ", str(value)).strip()[:48] for value in missing_texts[:2]]
        missing_summary = "；未覆盖：" + " / ".join(excerpts)
    if isinstance(coverage, (int, float)) and coverage < 0.7:
        issues.append({"severity": "P0", "category": "content", "message": f"可见原文逐块覆盖率仅 {coverage:.1%}{missing_summary}"})
    elif isinstance(coverage, (int, float)) and coverage < 0.98:
        issues.append({"severity": "P1", "category": "content", "message": f"可见原文逐块覆盖率为 {coverage:.1%}{missing_summary}"})
    picture = float(metrics.get("max_picture_coverage") or 0.0)
    if picture >= 0.85:
        issues.append({"severity": "P0", "category": "editability", "message": f"单个图片覆盖页面 {picture:.1%}，疑似整页截图"})
    elif picture >= 0.65:
        issues.append({"severity": "P1", "category": "editability", "message": f"单个图片覆盖页面 {picture:.1%}"})
    if int(metrics.get("out_of_bounds_count") or 0):
        issues.append({"severity": "P1", "category": "layout", "message": f"{metrics['out_of_bounds_count']} 个对象越出画布"})
    coarse = float(metrics.get("coarse_rgb_loss") or 0.0)
    ink = float(metrics.get("content_ink_loss") or 0.0)
    if coarse >= 0.1 or ink >= 0.08:
        issues.append({"severity": "P1", "category": "visual", "message": f"PowerPoint 视觉差异较大（RGB {coarse:.3f} / ink {ink:.3f}）"})
    elif coarse > 0.04 or ink > 0.03:
        issues.append({"severity": "P2", "category": "visual", "message": f"存在可见样式或几何差异（RGB {coarse:.3f} / ink {ink:.3f}）"})
    if any(item["severity"] == "P0" for item in issues):
        return "reject", issues
    if any(item["severity"] == "P1" for item in issues):
        return "needs_work", issues
    # P2 records reviewable decoration/style differences without rejecting a
    # business-usable editable page.  Final baseline promotion is still human.
    return "acceptable", issues


def _seconds(value: float | None) -> str:
    if value is None:
        return "—"
    return f"{value:.1f} 秒" if value < 60 else f"{int(value // 60)} 分 {value % 60:.1f} 秒"


def _report(
    page_dir: Path, case: PageCase, outcome: PageOutcome, run_meta: dict[str, Any],
    render_report: dict[str, Any], changed: list[str], render_elapsed: float | None,
) -> None:
    metrics = outcome.metrics
    lines = [
        f"# {case.page_id}", "",
        f"**结论：`{outcome.verdict}`**", "",
        "| 源图 | Microsoft PowerPoint 效果 |", "|---|---|",
        "| ![source](source.png) | " + ("![candidate](candidate.png)" if (page_dir / "candidate.png").is_file() else "未生成") + " |",
        "", "## 当前问题", "",
    ]
    if outcome.issues:
        for issue in outcome.issues:
            lines.append(f"- **{issue['severity']} · {issue['category']}**：{issue['message']}")
    else:
        lines.append("- 未发现自动诊断范围内的阻断问题；仍需人工查看视觉层级和局部样式。")
    if outcome.error:
        lines.append(f"- **执行错误**：{outcome.error}")
    lines += [
        "", "## 时间", "", "| 阶段 | 耗时 |", "|---|---:|",
        f"| 完整页面任务 | {_seconds(outcome.elapsed_sec)} |",
        f"| Codex 执行 | {_seconds(metrics.get('codex_elapsed_sec'))} |",
        f"| Microsoft PowerPoint 渲染 | {_seconds(render_elapsed)} |",
        "", "## PPTX 回读与视觉指标", "", "| 指标 | 值 |", "|---|---:|",
    ]
    for name in (
        "object_count", "text_shape_count", "native_shape_count", "table_count",
        "picture_count", "max_picture_coverage", "text_coverage", "missing_text_count",
        "excluded_expected_text_count",
        "coarse_rgb_loss", "content_ink_loss", "out_of_bounds_count",
        "codex_powerpoint_rendered",
    ):
        if name in metrics and metrics[name] is not None:
            lines.append(f"| `{name}` | `{metrics[name]}` |")
    missing_texts = metrics.get("missing_texts") or []
    if missing_texts:
        lines += ["", "### 未覆盖的可见原文", ""]
        for value in missing_texts:
            clean_value = re.sub(r"[\s\u000b]+", " ", str(value)).strip()
            lines.append(f"- {clean_value}")
    excluded_expected = metrics.get("excluded_expected_texts") or []
    if excluded_expected:
        lines += ["", "### 未纳入覆盖率的源对象", ""]
        for value in excluded_expected:
            lines.append(f"- `{value.get('reason', 'excluded')}`：{value.get('text', '')}")
    lines += ["", "## 关键决策", ""]
    if outcome.decisions:
        lines.extend(f"- {value.replace(chr(10), ' ')}" for value in outcome.decisions)
    else:
        lines.append("- 未记录。")
    lines += ["", "## 调用脚本与命令", ""]
    if outcome.skill_commands:
        lines.append("### `editppt` 确定性工具")
        lines.append("")
        for command in outcome.skill_commands:
            lines.append(
                f"- `{command.command}` → `{command.exit_code}`，{_seconds(command.elapsed_sec)}"
            )
        lines.append("")
        lines.append("### Codex shell 记录")
        lines.append("")
    if outcome.commands:
        for command in outcome.commands:
            safe_command = command.command.replace("`", "ˋ")
            lines.append(f"- `{safe_command}` → `{command.exit_code}`")
    else:
        lines.append("- 未记录。")
    if changed:
        lines += ["", "### Codex 新建或修改的过程文件", ""]
        lines.extend(f"- `{Path(value).name}`" for value in changed)
    skill = run_meta.get("skill") or {}
    renderer = render_report.get("renderer") or "microsoft-powerpoint"
    version = render_report.get("renderer_version") or "unknown"
    lines += [
        "", "## 运行身份", "", "| 项目 | 值 |", "|---|---|",
        f"| 输入 SHA-256 | `{_sha256(page_dir / 'source.png')}` |",
        f"| 输出 SHA-256 | `{_sha256(page_dir / 'candidate.pptx') if (page_dir / 'candidate.pptx').is_file() else '未生成'}` |",
        f"| Skill commit | `{skill.get('commit', '')}` |",
        f"| Skill dirty | `{skill.get('dirty', False)}` |",
        f"| Codex profile | `{run_meta.get('profile', '')}` |",
        f"| Model / effort | `{run_meta.get('model') or 'profile default'}` / `{run_meta.get('effort') or 'profile default'}` |",
        f"| Renderer | `{renderer}` / `{version}` |",
        "", "## 过程证据", "",
        "所有 OCR、manifest、临时 authoring 脚本、Codex JSONL、PowerPoint 报告、diff 和 heatmap 均位于 [`artifacts/`](artifacts/)。",
    ]
    (page_dir / "report.md").write_text("\n".join(lines) + "\n", encoding="utf-8")


def _run_page(
    args: argparse.Namespace, case: PageCase, run_dir: Path, run_meta: dict[str, Any],
) -> PageOutcome:
    page_dir = run_dir / "pages" / case.page_id
    artifacts = page_dir / "artifacts"
    work = artifacts / "authoring" / "work"
    codex_dir = artifacts / "codex"
    render_dir = artifacts / "render"
    compare_dir = artifacts / "compare"
    work.mkdir(parents=True, exist_ok=False)
    prepare_authoring_source(case.source, work)
    shutil.copy2(case.source, page_dir / "source.png")
    started = time.monotonic()
    outcome = PageOutcome(case.page_id)
    render_report: dict[str, Any] = {}
    render_elapsed: float | None = None
    changed: list[str] = []
    try:
        returncode, codex_elapsed, stderr = _execute_codex(
            args, Path(args.skill_root).resolve(), work, codex_dir
        )
        outcome.metrics["codex_elapsed_sec"] = round(codex_elapsed, 3)
        outcome.commands, outcome.decisions, changed = _events(codex_dir / "events.jsonl")
        outcome.skill_commands = _skill_trace(
            artifacts / "telemetry.jsonl",
            work / "editppt-events.jsonl",
        )
        if returncode != 0:
            raise BenchmarkError(f"Codex exited {returncode}: {stderr[-1000:]}")
        result = _read_json(work / "result.json")
        candidate = work / str(result.get("output_pptx") or "page.pptx")
        if result.get("status") != "ready" or not candidate.is_file() or candidate.stat().st_size == 0:
            raise BenchmarkError("Codex did not produce ready result.json and a non-empty page.pptx")
        shutil.copy2(candidate, page_dir / "candidate.pptx")
        candidate_sha = _sha256(page_dir / "candidate.pptx")
        bindings = [
            _read_json(path)
            for path in work.rglob("render-binding.json")
            if path.is_file()
        ]
        outcome.metrics["codex_powerpoint_rendered"] = any(
            value.get("authoritative") is True and value.get("input_sha256") == candidate_sha
            for value in bindings
        )
        rendered, render_report, render_elapsed, render_error = _render_powerpoint(
            page_dir / "candidate.pptx", render_dir, Path(args.skill_root).resolve(),
        )
        if rendered:
            shutil.copy2(rendered, page_dir / "candidate.png")
        metrics = _pptx_metrics(page_dir / "candidate.pptx", case.expected)
        if rendered:
            metrics.update(_visual_metrics(page_dir / "source.png", page_dir / "candidate.png", compare_dir))
        outcome.metrics.update(metrics)
        outcome.verdict, outcome.issues = _issues(
            outcome.metrics,
            render_error,
            outcome.commands,
            page_dir=page_dir,
            pages_root=run_dir / "pages",
        )
        if render_error:
            outcome.error = render_error
    except Exception as exc:  # noqa: BLE001 - evidence records all page failures
        outcome.error = str(exc)
        outcome.issues.append({"severity": "P0", "category": "execution", "message": str(exc)})
        outcome.verdict = "execution_failed"
    outcome.elapsed_sec = round(time.monotonic() - started, 3)
    telemetry = {
        "schema_version": 2, "page_id": case.page_id,
        "verdict": outcome.verdict, "elapsed_sec": outcome.elapsed_sec,
        "metrics": outcome.metrics, "issues": outcome.issues,
        "commands": [asdict(value) for value in outcome.commands],
        "skill_commands": [asdict(value) for value in outcome.skill_commands],
        "decisions": outcome.decisions, "error": outcome.error,
        "render_elapsed_sec": render_elapsed,
    }
    _write_json(artifacts / "telemetry.json", telemetry)
    _report(page_dir, case, outcome, run_meta, render_report, changed, render_elapsed)
    unexpected = {value.name for value in page_dir.iterdir()} - ROOT_FILES
    if unexpected:
        raise BenchmarkError(f"snapshot root contains process files: {sorted(unexpected)}")
    return outcome


def _summary(run_dir: Path, run_meta: dict[str, Any], outcomes: Iterable[PageOutcome]) -> None:
    values = list(outcomes)
    counts: dict[str, int] = {}
    for value in values:
        counts[value.verdict] = counts.get(value.verdict, 0) + 1
    issue_counts: dict[str, int] = {}
    for value in values:
        for issue in value.issues:
            key = f"{issue.get('severity')}:{issue.get('category')}"
            issue_counts[key] = issue_counts.get(key, 0) + 1
    lines = [
        f"# Image-to-Editable-PPT Benchmark：{run_meta['label']}", "",
        f"- Run ID：`{run_meta['run_id']}`",
        f"- Suite：`{run_meta['suite']}`",
        f"- Skill：`{run_meta['skill']['commit']}`",
        f"- 页面：{len(values)}",
        f"- 总耗时：{_seconds(sum(value.elapsed_sec for value in values))}",
        "", "## 结论分布", "", "| 结论 | 页数 |", "|---|---:|",
    ]
    lines.extend(f"| `{name}` | {count} |" for name, count in sorted(counts.items()))
    lines += ["", "## 页面", "", "| 页面 | 结论 | 耗时 | P0/P1/P2 |", "|---|---|---:|---|"]
    for value in values:
        severities = "/".join(
            str(sum(1 for issue in value.issues if issue.get("severity") == severity))
            for severity in ("P0", "P1", "P2")
        )
        lines.append(f"| [{value.page_id}](pages/{value.page_id}/report.md) | `{value.verdict}` | {_seconds(value.elapsed_sec)} | {severities} |")
    lines += ["", "## 问题分布", ""]
    if issue_counts:
        lines.extend(f"- `{name}`：{count}" for name, count in sorted(issue_counts.items()))
    else:
        lines.append("- 无自动诊断问题。")
    (run_dir / "SUMMARY.md").write_text("\n".join(lines) + "\n", encoding="utf-8")


def run_benchmark(args: argparse.Namespace) -> int:
    corpus = Path(args.corpus).expanduser().resolve()
    output = Path(args.out).expanduser().resolve()
    skill_root = Path(args.skill_root).expanduser().resolve()
    if not (skill_root / "SKILL.md").is_file():
        raise BenchmarkError(f"invalid Skill root: {skill_root}")
    cases = _filter_cases(_resolve_cases(corpus, args.suite), args.page)
    run_id = args.run_id or _run_id(args.label)
    run_dir = output / run_id
    run_dir.mkdir(parents=True, exist_ok=False)
    (run_dir / "pages").mkdir()
    repo = skill_root.parents[1]
    run_meta = {
        "schema_version": 2, "run_id": run_id, "label": args.label,
        "suite": args.suite, "started_at": _now(), "profile": args.profile,
        "page_filter": list(args.page or []),
        "model": args.model, "effort": args.effort,
        "skill": {
            "root": str(skill_root), "commit": _git_value(repo, "rev-parse", "HEAD"),
            "dirty": bool(_git_value(repo, "status", "--porcelain=v1")),
        },
        "corpus": str(corpus),
        "pages": [{"page_id": value.page_id, "source_sha256": _sha256(value.source)} for value in cases],
    }
    _write_json(run_dir / "run.json", run_meta)
    outcomes: list[PageOutcome] = []
    for index, case in enumerate(cases, start=1):
        print(f"[{index}/{len(cases)}] {case.page_id}", flush=True)
        outcome = _run_page(args, case, run_dir, run_meta)
        outcomes.append(outcome)
        print(f"  {outcome.verdict} in {outcome.elapsed_sec:.1f}s", flush=True)
    run_meta["completed_at"] = _now()
    run_meta["verdict_counts"] = {
        name: sum(1 for value in outcomes if value.verdict == name)
        for name in ("acceptable", "needs_work", "reject", "execution_failed")
    }
    _write_json(run_dir / "run.json", run_meta)
    _summary(run_dir, run_meta, outcomes)
    print(run_dir)
    return 0 if all(value.verdict == "acceptable" for value in outcomes) else 2


def verify_corpus(args: argparse.Namespace) -> int:
    corpus = Path(args.corpus).expanduser().resolve()
    cases = _filter_cases(_resolve_cases(corpus, args.suite), args.page)
    payload = {
        "ok": True, "suite": args.suite, "page_count": len(cases),
        "pages": [{"page_id": value.page_id, "sha256": _sha256(value.source)} for value in cases],
    }
    print(json.dumps(payload, ensure_ascii=False, indent=2))
    return 0


def build_parser() -> argparse.ArgumentParser:
    parser = argparse.ArgumentParser(description="Lean benchmark for one-Codex-task-per-page editable PPT rebuilding.")
    sub = parser.add_subparsers(dest="command", required=True)
    ingest = sub.add_parser("ingest-image")
    ingest.add_argument("--corpus", required=True)
    ingest.add_argument("--source", required=True)
    ingest.add_argument("--name", required=True)
    ingest.add_argument("--note", default="")
    ingest.add_argument("--text-hints", default="")
    ingest.set_defaults(func=ingest_image)
    for name in ("run", "verify"):
        command = sub.add_parser(name)
        command.add_argument("--corpus", required=True)
        command.add_argument("--suite", required=True)
        command.add_argument(
            "--page", action="append", default=[],
            help="run or verify only this exact page_id from the selected suite; repeatable",
        )
        if name == "run":
            command.add_argument("--out", required=True)
            command.add_argument("--label", required=True)
            command.add_argument("--run-id", default="")
            command.add_argument("--skill-root", required=True)
            command.add_argument("--profile", default="image-to-ppt-agent")
            command.add_argument("--model", default="")
            command.add_argument("--effort", default="")
            command.add_argument("--codex-bin", default="codex")
            command.set_defaults(func=run_benchmark)
        else:
            command.set_defaults(func=verify_corpus)
    return parser


def main(argv: list[str] | None = None) -> int:
    parser = build_parser()
    args = parser.parse_args(argv)
    try:
        return int(args.func(args) or 0)
    except BenchmarkError as exc:
        print(str(exc), file=sys.stderr)
        return 1
