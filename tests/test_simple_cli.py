import json
import os
import subprocess
import sys
import tempfile
import unittest
import zipfile
from pathlib import Path
from types import SimpleNamespace
from unittest import mock

from PIL import Image
from pptx import Presentation


ROOT = Path(__file__).resolve().parents[1]
RUNTIME = ROOT / "skills/image-to-editable-ppt/cli/editppt/runtime/main.py"
sys.path.insert(0, str(RUNTIME.parent))
import main as runtime_main  # noqa: E402
sys.path.pop(0)


def run_cli(*args: object, env=None) -> subprocess.CompletedProcess[str]:
    return subprocess.run(
        [sys.executable, str(RUNTIME), *[str(value) for value in args]],
        text=True,
        capture_output=True,
        env=env,
    )


def manifest(text: str) -> dict:
    return {
        "slide": {"width": 13.333, "height": 7.5, "background": "#FFFFFF"},
        "source": {"width_px": 1600, "height_px": 900},
        "shapes": [
            {
                "type": "rect",
                "box_px": [80, 180, 1440, 500],
                "fill": "#EAF3FF",
                "stroke": "#2375D8",
            }
        ],
        "images": [],
        "text_boxes": [
            {
                "text": text,
                "box_px": [120, 80, 1360, 90],
                "font_size": 28,
                "color": "#0A65B7",
            }
        ],
    }


class SimpleCliTest(unittest.TestCase):
    def test_help_exposes_only_simple_page_helpers(self):
        completed = run_cli("--help")
        self.assertEqual(0, completed.returncode, completed.stderr)
        for command in ("prepare", "inspect", "assets", "build", "text-fit", "render", "compare", "assemble", "doctor"):
            self.assertIn(command, completed.stdout)
        self.assertNotIn("dispatch", completed.stdout)
        self.assertNotIn("finalize", completed.stdout)

    def test_prepare_writes_ordered_source_pages_without_agent_state(self):
        with tempfile.TemporaryDirectory() as temporary:
            root = Path(temporary)
            first = root / "one.png"
            second = root / "two.png"
            Image.new("RGB", (160, 90), "white").save(first)
            Image.new("RGB", (160, 90), "blue").save(second)
            run_dir = root / "run"
            completed = run_cli("prepare", first, second, "--job-dir", run_dir)
            self.assertEqual(0, completed.returncode, completed.stderr)
            self.assertTrue((run_dir / "pages/page_001/source.png").is_file())
            self.assertTrue((run_dir / "pages/page_002/source.png").is_file())
            self.assertFalse((run_dir / "page_jobs.json").exists())
            payload = json.loads((run_dir / "deck_manifest.json").read_text())
            assert "agent_status" not in json.dumps(payload)
            assert "validation" not in json.dumps(payload)
            assert payload["pages"][0]["result"].endswith("result.json")

    def test_prepare_uses_stable_authoring_space_and_crop_restores_original_pixels(self):
        with tempfile.TemporaryDirectory() as temporary:
            root = Path(temporary)
            source = root / "wide.png"
            image = Image.new("RGB", (4096, 2304), "white")
            for x in range(2000, 2400):
                for y in range(800, 1200):
                    image.putpixel((x, y), (0, 133, 208))
            image.save(source)
            run_dir = root / "run"

            prepared = run_cli("prepare", source, "--job-dir", run_dir)
            self.assertEqual(0, prepared.returncode, prepared.stderr)
            page = run_dir / "pages/page_001"
            with Image.open(page / "source.png") as authoring:
                self.assertEqual((2048, 1152), authoring.size)
            mapping = json.loads((page / ".editppt/source-map.json").read_text())
            self.assertEqual([4096, 2304], mapping["original_size_px"])

            cropped = page / "asset.png"
            result = run_cli(
                "assets", "crop", "--input", page / "source.png", "--out", cropped,
                "--left", 1000, "--top", 400, "--right", 1200, "--bottom", 600,
            )
            self.assertEqual(0, result.returncode, result.stderr)
            payload = json.loads(result.stdout)
            self.assertEqual([1000, 400, 1200, 600], payload["authoring_box_px"])
            self.assertEqual([2000, 800, 2400, 1200], payload["box_px"])
            with Image.open(cropped) as asset:
                self.assertEqual((400, 400), asset.size)

    def test_build_draft_preview_and_inspect_native_table(self):
        with tempfile.TemporaryDirectory() as temporary:
            root = Path(temporary)
            page = root / "page"
            page.mkdir()
            Image.new("RGB", (1600, 900), "white").save(page / "source.png")
            value = manifest("第一页")
            value["tables"] = [{
                "box_px": [100, 250, 1400, 300],
                "rows": [["列一", "列二"], ["内容 A", "内容 B"]],
                "header_fill": "#0A65B7",
                "header_color": "#FFFFFF",
            }]
            (page / "manifest.json").write_text(json.dumps(value, ensure_ascii=False), encoding="utf-8")
            built = run_cli("build", page, "--draft-preview", "draft.png")
            self.assertEqual(0, built.returncode, built.stderr)
            self.assertTrue((page / "page.pptx").is_file())
            self.assertTrue((page / "draft.png").is_file())
            self.assertFalse((page / "preview.png").exists())
            inspected = run_cli("inspect", "pptx", page)
            self.assertEqual(0, inspected.returncode, inspected.stderr)
            payload = json.loads(inspected.stdout)
            self.assertEqual(1, payload["summary"]["table_count"])
            with zipfile.ZipFile(page / "page.pptx") as package:
                self.assertIn("ppt/presProps.xml", package.namelist())
                self.assertIn("ppt/viewProps.xml", package.namelist())
                self.assertIn("docProps/thumbnail.jpeg", package.namelist())

    def test_assemble_passes_independent_page_pptx_in_order(self):
        with tempfile.TemporaryDirectory() as temporary:
            root = Path(temporary)
            pages = []
            for index, title in enumerate(("第一页", "第二页"), start=1):
                page = root / f"page-{index}"
                page.mkdir()
                Image.new("RGB", (30, 20), (index * 40, 100, 180)).save(page / "asset.png")
                value = manifest(title)
                value["images"] = [{"path": "asset.png", "box_px": [1400, 30, 100, 60]}]
                value["tables"] = [{"box_px": [100, 300, 600, 200], "rows": [["页", "值"], [title, index]]}]
                (page / "manifest.json").write_text(json.dumps(value, ensure_ascii=False), encoding="utf-8")
                self.assertEqual(0, run_cli("build", page).returncode)
                pages.append(page)
            output = root / "deck.pptx"
            assembled = run_cli("assemble", *pages, "--out", output, "--evidence-dir", root / "evidence")
            self.assertEqual(0, assembled.returncode, assembled.stderr)
            deck = Presentation(str(output))
            self.assertEqual(2, len(deck.slides))
            self.assertIn("第一页", "\n".join(shape.text for shape in deck.slides[0].shapes if hasattr(shape, "text")))
            self.assertIn("第二页", "\n".join(shape.text for shape in deck.slides[1].shapes if hasattr(shape, "text")))
            for slide in deck.slides:
                self.assertTrue(any(getattr(shape, "has_table", False) for shape in slide.shapes))
                self.assertTrue(any(str(shape.shape_type) == "PICTURE (13)" for shape in slide.shapes))

    def test_doctor_reports_contract_version(self):
        completed = run_cli("doctor", "--json")
        payload = json.loads(completed.stdout)
        self.assertEqual("benchmark-driven-page-v3", payload["skill"]["contract_version"])
        powerpoint_ready = payload["editppt_checks"]["powerpoint"]["available"]
        renderer_ready = payload["editppt_checks"]["powerpoint_renderer"]["available"]
        self.assertEqual(0 if powerpoint_ready and renderer_ready else 1, completed.returncode)

    def test_inspect_uses_configured_content_aware_ocr_without_exposing_token(self):
        with tempfile.TemporaryDirectory() as temporary:
            page = Path(temporary)
            Image.new("RGB", (160, 90), "white").save(page / "source.png")
            calls: list[tuple[str, dict[str, str] | None]] = []

            def fake_script(name: str, *argv: object, capture: bool = False, env=None):
                calls.append((name, env))
                (page / "text_hints.json").write_text(
                    json.dumps({"backend": "paddleocr-vl", "lines": []}),
                    encoding="utf-8",
                )
                return subprocess.CompletedProcess([], 0, "", "")

            args = SimpleNamespace(
                page_dir=page,
                source="source.png",
                out="text_hints.json",
                overlay="text_hints.png",
            )
            with mock.patch.object(runtime_main, "_configured_paddle_token", return_value="secret-value"), mock.patch.object(
                runtime_main, "_script", side_effect=fake_script
            ):
                self.assertEqual(0, runtime_main.cmd_inspect_text(args))

            self.assertEqual("paddle_text_hints.py", calls[0][0])
            self.assertEqual("secret-value", calls[0][1]["PADDLE_OCR_TOKEN"])

    def test_trace_records_actual_command_exit_and_duration(self):
        with tempfile.TemporaryDirectory() as temporary:
            trace = Path(temporary) / "trace.jsonl"
            environment = os.environ.copy()
            environment["EDITPPT_TRACE_FILE"] = str(trace)
            completed = run_cli("text-fit", "--text", "单行标题", "--width-px", 600, "--height-px", 80, "--single-line", env=environment)
            self.assertEqual(0, completed.returncode, completed.stderr)
            event = json.loads(trace.read_text(encoding="utf-8"))
            self.assertEqual("text-fit", event["command"])
            self.assertEqual(0, event["exit_code"])
            self.assertGreaterEqual(event["elapsed_sec"], 0)

    def test_source_crop_and_flat_background_separation_are_audited(self):
        with tempfile.TemporaryDirectory() as temporary:
            root = Path(temporary)
            source = root / "source.png"
            image = Image.new("RGB", (200, 100), "white")
            for x in range(60, 140):
                for y in range(25, 75):
                    image.putpixel((x, y), (0, 133, 208))
            image.save(source)
            cropped = root / "crop.png"
            result = run_cli(
                "assets", "crop", "--input", source, "--out", cropped,
                "--left", 50, "--top", 20, "--right", 150, "--bottom", 80,
            )
            self.assertEqual(0, result.returncode, result.stderr)
            crop_report = json.loads(result.stdout)
            self.assertFalse(crop_report["near_full_page_risk"])
            separated = root / "separated.png"
            result = run_cli("assets", "separate", "--input", cropped, "--out", separated)
            self.assertEqual(0, result.returncode, result.stderr)
            separation = json.loads(result.stdout)
            self.assertLess(separation["put_back_mae"], 0.01)
            with Image.open(separated) as separated_image:
                self.assertEqual("RGBA", separated_image.mode)

    def test_layout_and_structure_use_source_coordinates_without_object_ids(self):
        with tempfile.TemporaryDirectory() as temporary:
            page = Path(temporary)
            image = Image.new("RGB", (400, 225), "white")
            for x in range(50, 350):
                image.putpixel((x, 60), (0, 133, 208))
                image.putpixel((x, 160), (0, 133, 208))
            for y in range(60, 161):
                image.putpixel((50, y), (0, 133, 208))
                image.putpixel((349, y), (0, 133, 208))
            image.save(page / "source.png")
            layout = run_cli("inspect", "layout", page)
            structure = run_cli("inspect", "structure", page)
            self.assertEqual(0, layout.returncode, layout.stderr)
            self.assertEqual(0, structure.returncode, structure.stderr)
            payload = json.loads(structure.stdout)
            self.assertNotIn("parent_id", json.dumps(payload))
            self.assertGreater(len(payload["horizontal_segments"]), 0)


if __name__ == "__main__":
    unittest.main()
