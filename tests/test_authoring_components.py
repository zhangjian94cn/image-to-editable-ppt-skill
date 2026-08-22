from __future__ import annotations

import json
import subprocess
import sys
import tempfile
import zipfile
from pathlib import Path

from PIL import Image
from pptx import Presentation
import pytest


ROOT = Path(__file__).resolve().parents[1]
CLI_ROOT = ROOT / "skills/image-to-editable-ppt/cli"
RUNTIME = CLI_ROOT / "editppt/runtime/main.py"
sys.path.insert(0, str(CLI_ROOT))
from editppt.authoring import SlideManifest  # noqa: E402


def run_cli(*args: object) -> subprocess.CompletedProcess[str]:
    return subprocess.run(
        [sys.executable, str(RUNTIME), *[str(value) for value in args]],
        text=True,
        capture_output=True,
    )


def test_components_build_native_connector_and_rich_table():
    with tempfile.TemporaryDirectory() as temporary:
        page = Path(temporary)
        Image.new("RGB", (1600, 900), "white").save(page / "source.png")
        manifest = SlideManifest(1600, 900)
        manifest.add_section_band([80, 60, 1440, 64], "原生组件", fill="#0A65B7")
        manifest.add_shape(
            [80, 140, 1440, 40],
            gradient={
                "angle": 0,
                "stops": [
                    {"position": 0, "color": "#EAF3FF"},
                    {"position": 1, "color": "#FFFFFF"},
                ],
            },
        )
        manifest.add_timeline_stage([120, 220, 280, 48], number="1", period="第1-2周", heading="完成闭环")
        manifest.add_connector([400, 244], [620, 244], end_arrow="triangle")
        manifest.add_table(
            [120, 420, 1360, 260],
            [
                ["阶段", "说明"],
                ["准备", {"runs": [{"text": "可编辑", "bold": True, "color": "#D92D20"}, {"text": "单元格"}]}],
            ],
            column_widths=[1, 3],
        )
        manifest.write(page / "manifest.json")

        built = run_cli("build", page)
        assert built.returncode == 0, built.stderr
        build_payload = json.loads(built.stdout)
        assert Path(build_payload["build_report"]).is_file()
        inspected = run_cli("inspect", "pptx", page)
        assert inspected.returncode == 0, inspected.stderr
        payload = json.loads(inspected.stdout)
        assert payload["summary"]["connector_count"] == 1
        assert payload["summary"]["table_count"] == 1
        with zipfile.ZipFile(page / "page.pptx") as package:
            slide_xml = package.read("ppt/slides/slide1.xml").decode("utf-8")
        assert "<p:cxnSp>" in slide_xml
        assert '<a:tailEnd type="triangle"/>' in slide_xml
        assert '<a:gradFill rotWithShape="1">' in slide_xml
        assert "可编辑" in slide_xml and "单元格" in slide_xml


def test_builder_preserves_powerpoint_private_font_and_protects_single_line_title():
    with tempfile.TemporaryDirectory() as temporary:
        page = Path(temporary)
        Image.new("RGB", (1600, 900), "white").save(page / "source.png")
        value = SlideManifest(1600, 900)
        value.add_text(
            [53, 34, 1040, 53],
            "各阶段里程碑 | 看闭环、看资产，而不是只看功能数量",
            font="Microsoft YaHei",
            font_size=30.5,
            font_size_source="measured",
            bold=True,
        )
        value.write(page / "manifest.json")
        built = run_cli("build", page)
        assert built.returncode == 0, built.stderr
        build_payload = json.loads(built.stdout)
        resolution = build_payload["font_resolution"][0]
        if Path("/Applications/Microsoft PowerPoint.app/Contents/Resources/DFonts/msyh.ttc").is_file():
            assert not build_payload["font_substitutions"]
            assert resolution["resolved"] == "Microsoft YaHei"
            assert resolution["provider"] == "powerpoint-dfonts"
        elif resolution["path"]:
            assert build_payload["font_substitutions"]
            assert resolution["resolved"] != "Microsoft YaHei"
        else:
            assert not build_payload["font_substitutions"]
            assert resolution["provider"] == ""
        assert build_payload["text_adjustments"]
        deck = Presentation(str(page / "page.pptx"))
        shape = next(shape for shape in deck.slides[0].shapes if getattr(shape, "text", ""))
        run = shape.text_frame.paragraphs[0].runs[0]
        if resolution["resolved"] == "Microsoft YaHei":
            assert run.font.name == "Microsoft YaHei"
        else:
            assert run.font.name != "Microsoft YaHei"
        assert run.font.size.pt < 30.5


def test_layered_header_uses_named_layers_and_masks_accent_behind_band():
    with tempfile.TemporaryDirectory() as temporary:
        page = Path(temporary)
        Image.new("RGB", (1600, 900), "white").save(page / "source.png")
        value = SlideManifest(1600, 900)
        parts = value.add_layered_header(
            [10, 45, 420, 300],
            [10, 45, 420, 58],
            [90, 91, 260, 18],
            "整体使用分析",
            title_font_size_px=32,
        )
        assert parts["container"]["layer"] == "container"
        assert parts["accent"]["layer"] == "decoration_behind"
        assert parts["band"]["layer"] == "band"
        assert parts["title"]["layer"] == "text"
        value.write(page / "manifest.json")

        built = run_cli("build", page)

        assert built.returncode == 0, built.stderr
        payload = json.loads(built.stdout)
        report = json.loads(Path(payload["layer_report"]).read_text())
        ordered = [(item["layer"], item["z_index"]) for item in report["objects"]]
        assert ordered == [
            ("container", 100.0),
            ("decoration_behind", 120.0),
            ("band", 140.0),
            ("text", 300.0),
        ]
        assert not report["conflicts"]


def test_editable_chart_components_build_native_shapes_and_consistent_caption_roles():
    with tempfile.TemporaryDirectory() as temporary:
        page = Path(temporary)
        Image.new("RGB", (1600, 900), "white").save(page / "source.png")
        value = SlideManifest(1600, 900)
        bars = value.add_editable_bar_chart(
            [80, 80, 650, 330],
            ["泰州", "镇江", "徐州"],
            [19, 17, 13],
            title="各地市TPD",
            maximum=20,
            grid_values=[0, 5, 10, 15, 20],
        )
        line = value.add_editable_line_chart(
            [820, 80, 650, 330],
            ["02月", "03月", "04月", "05月"],
            [19.3, 23.2, 44.2, 46.6],
            title="Token 整体使用情况",
            minimum=10,
            maximum=55,
            grid_values=[10, 25, 40, 55],
            marker_color="#F05A67",
            layer="content",
        )
        value.write(page / "manifest.json")

        assert len(bars["bars"]) == 3
        assert len(line["line"]) == 5
        assert all(item["layer"] == "decoration_behind" for item in bars["grid"] + line["grid"])
        caption_sizes = {
            item["font_size"]
            for item in value.text_boxes
            if item["text_role"] == "caption" and "font_size" in item
        }
        assert len(caption_sizes) == 1

        built = run_cli("build", page)
        assert built.returncode == 0, built.stderr
        payload = json.loads(built.stdout)
        report = json.loads(Path(payload["build_report"]).read_text())
        assert report["objects"]["images"] == 0
        assert report["objects"]["shapes"] >= 17
        assert not [item for item in report["role_size_deviations"] if not item["measured_override"]]
        with zipfile.ZipFile(page / "page.pptx") as package:
            slide_xml = package.read("ppt/slides/slide1.xml").decode("utf-8")
        assert "各地市TPD" in slide_xml
        assert "Token 整体使用情况" in slide_xml
        assert '<a:path w="21600" h="21600" fill="none">' in slide_xml


def test_editable_chart_components_reject_invalid_data_contracts():
    page = SlideManifest(1600, 900)
    with pytest.raises(ValueError, match="equal length"):
        page.add_editable_bar_chart([10, 10, 400, 250], ["A"], [1, 2])
    with pytest.raises(ValueError, match="cover every value"):
        page.add_editable_bar_chart([10, 10, 400, 250], ["A"], [2], maximum=1)
    with pytest.raises(ValueError, match="at least two"):
        page.add_editable_line_chart([10, 10, 400, 250], ["A"], [1])


def test_builder_warns_when_overlapping_text_duplicates_a_visible_prefix():
    with tempfile.TemporaryDirectory() as temporary:
        page = Path(temporary)
        Image.new("RGB", (1600, 900), "white").save(page / "source.png")
        value = SlideManifest(1600, 900)
        value.add_text([100, 100, 600, 60], "分析结论：Token消耗增速放缓", font_size_px=20)
        value.add_text([100, 100, 130, 60], "分析结论：", font_size_px=20, color="#0085D0")
        value.write(page / "manifest.json")

        built = run_cli("build", page)

        assert built.returncode == 0, built.stderr
        payload = json.loads(built.stdout)
        report = json.loads(Path(payload["layer_report"]).read_text())
        assert report["duplicate_text_overlaps"][0]["duplicated_excerpt"] == "分析结论："
        assert any("duplicate the same visible content" in warning for warning in payload["warnings"])


def test_explicit_z_index_wins_over_named_layer_with_visible_warning():
    with tempfile.TemporaryDirectory() as temporary:
        page = Path(temporary)
        Image.new("RGB", (1600, 900), "white").save(page / "source.png")
        value = SlideManifest(1600, 900)
        value.add_shape([100, 100, 300, 100], fill="#0087CF", layer="band", z_index=500)
        value.write(page / "manifest.json")

        built = run_cli("build", page)

        assert built.returncode == 0, built.stderr
        payload = json.loads(built.stdout)
        report = json.loads(Path(payload["layer_report"]).read_text())
        assert report["objects"][0]["z_index"] == 500.0
        assert report["conflicts"][0]["conflict"]["layer_z"] == 140
        assert any("z_index won" in warning for warning in payload["warnings"])


def test_typography_roles_apply_page_scaled_defaults_and_warn_on_core_shrink():
    with tempfile.TemporaryDirectory() as temporary:
        page = Path(temporary)
        Image.new("RGB", (3200, 1800), "white").save(page / "source.png")
        value = SlideManifest(3200, 1800)
        value.add_text([100, 100, 100, 30], "放不下的页面主标题", text_role="slide_title")
        value.write(page / "manifest.json")

        built = run_cli("build", page)

        assert built.returncode == 0, built.stderr
        payload = json.loads(built.stdout)
        adjustment = payload["typography_adjustments"][0]
        assert adjustment["text_role"] == "slide_title"
        assert adjustment["font_size_source"] == "role_fallback"
        assert any("role shrink threshold" in warning for warning in payload["warnings"])
        shrink = payload["role_shrink_warnings"][0]
        assert shrink["text_excerpt"] == "放不下的页面主标题"
        assert shrink["geometry_diagnostic"]["limiting_dimension"] in {"width", "height", "unknown"}
        if shrink["geometry_diagnostic"]["limiting_dimension"] == "unknown":
            assert shrink["geometry_diagnostic"]["warning"]
        else:
            assert len(shrink["geometry_diagnostic"]["required_content_px"]) == 2


def test_role_deviation_compares_requested_size_not_fit_result():
    with tempfile.TemporaryDirectory() as temporary:
        page = Path(temporary)
        Image.new("RGB", (1600, 900), "white").save(page / "source.png")
        value = SlideManifest(1600, 900)
        value.add_text([20, 20, 1200, 60], "宽标题", font_size_px=32, text_role="section_title")
        value.add_text([20, 100, 120, 20], "同字号但框很小", font_size_px=32, text_role="section_title")
        value.write(page / "manifest.json")

        built = run_cli("build", page)

        assert built.returncode == 0, built.stderr
        payload = json.loads(built.stdout)
        assert payload["role_shrink_warnings"]
        assert payload["role_size_deviations"] == []


def test_measured_same_role_size_override_is_auditable_not_a_warning():
    with tempfile.TemporaryDirectory() as temporary:
        page = Path(temporary)
        Image.new("RGB", (1600, 900), "white").save(page / "source.png")
        value = SlideManifest(1600, 900)
        value.add_text([20, 20, 800, 60], "标题一", font_size_px=32, text_role="section_title")
        value.add_text(
            [20, 100, 800, 60],
            "标题二",
            font_size_px=28,
            text_role="section_title",
            font_size_source="measured",
        )
        value.write(page / "manifest.json")

        built = run_cli("build", page)

        assert built.returncode == 0, built.stderr
        payload = json.loads(built.stdout)
        assert payload["role_size_deviations"]
        assert all("role median" not in warning for warning in payload["warnings"])


def test_role_range_diagnostic_catches_uniformly_too_small_unmeasured_text():
    with tempfile.TemporaryDirectory() as temporary:
        page = Path(temporary)
        Image.new("RGB", (1600, 900), "white").save(page / "source.png")
        value = SlideManifest(1600, 900)
        value.add_text([20, 20, 300, 30], "标签一", font_size_px=9, text_role="caption")
        value.add_text([340, 20, 300, 30], "标签二", font_size_px=9, text_role="caption")
        value.add_text(
            [660, 20, 300, 30], "源图实测小标签", font_size_px=9,
            text_role="caption", font_size_source="measured",
        )
        value.write(page / "manifest.json")

        built = run_cli("build", page)

        assert built.returncode == 0, built.stderr
        payload = json.loads(built.stdout)
        report = json.loads(Path(payload["build_report"]).read_text())
        assert len(report["role_range_deviations"]) == 3
        assert sum(not item["measured_override"] for item in report["role_range_deviations"]) == 2
        assert any("semantic role size range" in warning for warning in payload["warnings"])


def test_builder_accepts_common_aliases_and_source_pixel_font_sizes():
    with tempfile.TemporaryDirectory() as temporary:
        page = Path(temporary)
        Image.new("RGB", (1600, 900), "white").save(page / "source.png")
        value = SlideManifest(1600, 900)
        value.add_shape([100, 100, 300, 80], "roundRect", fill="#EAF3FF", radius=12)
        value.add_connector([400, 140], [600, 140], stroke="#0A65B7", stroke_width=2)
        value.add_text([120, 110, 250, 50], "", font_size_px=30, valign="mid")
        value.add_table([100, 250, 600, 160], [["表头"], [{"text": "内容", "font_size_px": 20}]], font_size_px=24)
        value.write(page / "manifest.json")

        built = run_cli("build", page)
        assert built.returncode == 0, built.stderr
        with zipfile.ZipFile(page / "page.pptx") as package:
            slide_xml = package.read("ppt/slides/slide1.xml").decode("utf-8")
        assert 'anchor="ctr"' in slide_xml
        assert '<a:gd name="adj" fmla="val 15000"/>' in slide_xml
        assert 'w="25400"' in slide_xml
        assert 'sz="1800"' in slide_xml
        assert 'sz="1440"' in slide_xml
        assert 'sz="1200"' in slide_xml


def test_polygon_helper_uses_source_space_points():
    value = SlideManifest(1600, 900)
    item = value.add_polygon(
        [[10, 20], [110, 20], [80, 90], [20, 70]],
        fill="#DDEEFF",
        stroke="#0A65B7",
    )

    assert item["box_px"] == [10.0, 20.0, 100.0, 70.0]
    assert item["polygon_px"] == [[10.0, 20.0], [110.0, 20.0], [80.0, 90.0], [20.0, 70.0]]


def test_shape_helper_accepts_kind_first_and_open_polyline_points():
    value = SlideManifest(1600, 900)
    rect = value.add_shape("rect", box_px=[20, 30, 100, 50], fill="#EAF3FF")
    positional_rect = value.add_shape("roundRect", [140, 30, 100, 50], fill="#EAF3FF")
    path = value.add_shape(
        "line",
        points_px=[[100, 100], [160, 100], [160, 180], [240, 180]],
        stroke="#0A65B7",
        stroke_width=2,
        end_arrow="triangle",
    )

    assert rect["box_px"] == [20.0, 30.0, 100.0, 50.0]
    assert positional_rect["box_px"] == [140.0, 30.0, 100.0, 50.0]
    assert path["type"] == "polyline"
    assert path["polyline_px"][-1] == [240.0, 180.0]


def test_polyline_builds_an_open_editable_freeform_path():
    with tempfile.TemporaryDirectory() as temporary:
        page = Path(temporary)
        Image.new("RGB", (1600, 900), "white").save(page / "source.png")
        value = SlideManifest(1600, 900)
        value.add_polyline([[100, 100], [400, 100], [400, 250]], end_arrow="triangle")
        value.write(page / "manifest.json")

        built = run_cli("build", page)

        assert built.returncode == 0, built.stderr
        with zipfile.ZipFile(page / "page.pptx") as package:
            slide_xml = package.read("ppt/slides/slide1.xml").decode("utf-8")
        assert '<a:path w="21600" h="21600" fill="none">' in slide_xml
        assert "<a:close/>" not in slide_xml
        assert '<a:tailEnd type="triangle"/>' in slide_xml


def test_builder_surfaces_severe_automatic_text_shrink():
    with tempfile.TemporaryDirectory() as temporary:
        page = Path(temporary)
        Image.new("RGB", (1600, 900), "white").save(page / "source.png")
        value = SlideManifest(1600, 900)
        value.add_text([10, 10, 120, 20], "这是一个明显放不下的很长标题", font_size=32)
        value.write(page / "manifest.json")

        built = run_cli("build", page)
        assert built.returncode == 0, built.stderr
        payload = json.loads(built.stdout)
        assert payload["severe_text_adjustments"]
        assert payload["warnings"]
