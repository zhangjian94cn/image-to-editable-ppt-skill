from __future__ import annotations

import json
from argparse import Namespace
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.util import Inches

from benchmark_runner.runner import (
    CommandEvidence,
    ROOT_FILES,
    _codex_command,
    _issues,
    _filter_cases,
    _pptx_metrics,
    _resolve_cases,
    _skill_trace,
    _visual_metrics,
)


def _corpus(tmp_path: Path) -> Path:
    root = tmp_path / "corpus"
    page = root / "decks/case/slides/001"
    page.mkdir(parents=True)
    Image.new("RGB", (1600, 900), "white").save(page / "input.png")
    (page / "expected.json").write_text(json.dumps({"objects": [{"kind": "text", "text": "标题", "visible": True}]}))
    (root / "decks/case/deck.json").write_text(json.dumps({"slides": [{"number": 1, "path": "slides/001"}]}))
    (root / "suites.json").write_text(json.dumps({"page_suites": {"round-0": [{"deck_id": "case", "slide": 1}]}}))
    return root


def test_page_suite_resolves_one_frozen_page(tmp_path: Path):
    cases = _resolve_cases(_corpus(tmp_path), "round-0")
    assert [value.page_id for value in cases] == ["case-p001"]


def test_page_filter_keeps_exact_suite_page(tmp_path: Path):
    cases = _resolve_cases(_corpus(tmp_path), "round-0")
    selected = _filter_cases(cases, ["case-p001"])
    assert selected == cases


def test_pptx_readback_flags_full_page_picture_and_text(tmp_path: Path):
    pptx = tmp_path / "candidate.pptx"
    picture = tmp_path / "picture.png"
    Image.new("RGB", (1600, 900), "white").save(picture)
    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    slide.shapes.add_picture(str(picture), 0, 0, presentation.slide_width, presentation.slide_height)
    slide.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(1)).text = "标题"
    presentation.save(pptx)
    expected = tmp_path / "expected.json"
    expected.write_text(json.dumps({"objects": [{"kind": "text", "text": "标题", "visible": True}]}))
    metrics = _pptx_metrics(pptx, expected)
    assert metrics["text_coverage"] == 1.0
    assert metrics["max_picture_coverage"] > 0.99
    verdict, issues = _issues({**metrics, "coarse_rgb_loss": 0.0, "content_ink_loss": 0.0})
    assert verdict == "reject"
    assert any(value["category"] == "editability" for value in issues)


def test_text_coverage_ignores_dynamic_placeholders_and_text_hidden_by_later_picture(tmp_path: Path):
    pptx = tmp_path / "candidate.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1)).text = "标题 | 可编辑页面"
    presentation.save(pptx)
    expected = tmp_path / "expected.json"
    expected.write_text(json.dumps({"objects": [
        {"kind": "text", "text": "‹#›", "box_px": [900, 800, 60, 30]},
        {"kind": "text", "text": "标题丨可编辑页面", "box_px": [100, 100, 500, 80]},
        {"kind": "text", "text": "应用截图", "box_px": [700, 300, 60, 180]},
        {"kind": "picture", "box_px": [650, 250, 200, 350]},
    ]}, ensure_ascii=False))

    metrics = _pptx_metrics(pptx, expected)

    assert metrics["text_coverage"] == 1.0
    assert metrics["expected_text_count"] == 1
    assert metrics["excluded_expected_text_count"] == 2
    assert {value["reason"] for value in metrics["excluded_expected_texts"]} == {
        "dynamic_placeholder", "occluded_by_later_picture",
    }


def test_text_coverage_reports_exact_visible_missing_text(tmp_path: Path):
    pptx = tmp_path / "candidate.pptx"
    presentation = Presentation()
    presentation.slides.add_slide(presentation.slide_layouts[6])
    presentation.save(pptx)
    expected = tmp_path / "expected.json"
    expected.write_text(json.dumps({"objects": [
        {"kind": "text", "text": "不能遗漏的版权句", "box_px": [10, 10, 300, 30]},
    ]}, ensure_ascii=False))

    metrics = _pptx_metrics(pptx, expected)

    assert metrics["text_coverage"] == 0.0
    assert metrics["missing_texts"] == ["不能遗漏的版权句"]


def test_visual_ink_metric_ignores_light_background_style_changes(tmp_path: Path):
    source = tmp_path / "source.png"
    candidate = tmp_path / "candidate.png"
    Image.new("RGB", (200, 100), (225, 225, 225)).save(source)
    Image.new("RGB", (200, 100), "white").save(candidate)

    metrics = _visual_metrics(source, candidate, tmp_path / "compare")

    assert metrics["content_ink_loss"] == 0.0
    assert metrics["coarse_rgb_loss"] > 0.0


def test_p2_visual_difference_remains_acceptable_for_human_review():
    metrics = {
        "codex_powerpoint_rendered": True,
        "text_coverage": 1.0,
        "max_picture_coverage": 0.1,
        "out_of_bounds_count": 0,
        "coarse_rgb_loss": 0.05,
        "content_ink_loss": 0.02,
    }

    verdict, issues = _issues(metrics)

    assert verdict == "acceptable"
    assert [value["severity"] for value in issues] == ["P2"]


def test_direct_powerpoint_automation_is_rejected():
    metrics = {
        "codex_powerpoint_rendered": True,
        "text_coverage": 1.0,
        "max_picture_coverage": 0.1,
        "out_of_bounds_count": 0,
        "coarse_rgb_loss": 0.0,
        "content_ink_loss": 0.0,
    }
    verdict, issues = _issues(
        metrics,
        commands=[CommandEvidence("osascript -e 'tell application Microsoft PowerPoint to activate'", 0)],
    )
    assert verdict == "reject"
    assert any(value["category"] == "execution" for value in issues)


def test_parent_or_sibling_page_scan_is_reported():
    metrics = {
        "codex_powerpoint_rendered": True,
        "text_coverage": 1.0,
        "max_picture_coverage": 0.1,
        "out_of_bounds_count": 0,
        "coarse_rgb_loss": 0.0,
        "content_ink_loss": 0.0,
    }
    root = Path("/tmp/snapshot/pages")
    page = root / "case-p001"
    verdict, issues = _issues(
        metrics,
        commands=[CommandEvidence("find .. -name manifest.json", 0)],
        page_dir=page,
        pages_root=root,
    )
    assert verdict == "needs_work"
    assert any("单页独立性" in value["message"] for value in issues)


def test_historical_memory_read_is_reported():
    metrics = {
        "codex_powerpoint_rendered": True,
        "text_coverage": 1.0,
        "max_picture_coverage": 0.1,
        "out_of_bounds_count": 0,
        "coarse_rgb_loss": 0.0,
        "content_ink_loss": 0.0,
    }
    verdict, issues = _issues(
        metrics,
        commands=[CommandEvidence("rg image-to-ppt /Users/test/.codex/memories/MEMORY.md", 0)],
    )
    assert verdict == "needs_work"
    assert any("历史 memory" in value["message"] for value in issues)


def test_snapshot_root_contract_is_review_first():
    assert ROOT_FILES == {"source.png", "candidate.pptx", "candidate.png", "report.md", "artifacts"}


def test_benchmark_codex_has_host_access_for_powerpoint(tmp_path: Path):
    args = Namespace(codex_bin="/usr/bin/true", profile="image-to-ppt-agent", model="", effort="")
    command = _codex_command(args, tmp_path, tmp_path / "source.png")
    assert command[command.index("-s") + 1] == "danger-full-access"


def test_skill_trace_reads_runner_and_page_local_files(tmp_path: Path):
    runner_trace = tmp_path / "telemetry.jsonl"
    page_trace = tmp_path / "editppt-events.jsonl"
    runner_trace.write_text(json.dumps({"command": "inspect", "subcommand": "text", "exit_code": 0, "elapsed_sec": 1.25}) + "\n")
    page_trace.write_text(json.dumps({"command": "build", "exit_code": 0, "elapsed_sec": 0.5}) + "\n")
    evidence = _skill_trace(runner_trace, page_trace)
    assert [value.command for value in evidence] == ["editppt inspect text", "editppt build"]
