#!/usr/bin/env python3
import argparse
import html
import json
import math
import re
import subprocess
import sys
import tempfile
import zipfile
from copy import deepcopy
from pathlib import Path
from xml.etree import ElementTree as ET

from PIL import ImageFont

from text_fit_tool import find_font


EMU_PER_INCH = 914400
REL_NS = "http://schemas.openxmlformats.org/package/2006/relationships"
ASPECT_16_9 = 16 / 9
ASPECT_TOLERANCE = 0.03
DEFAULT_TEXT_FIT_SAFETY = 0.9
DEFAULT_TEXT_LINE_HEIGHT = 1.22
DEFAULT_MIN_FONT_SIZE = 4.0
TEXT_ALIGNMENTS = {
    "left": ("l", "left"),
    "center": ("ctr", "center"),
    "right": ("r", "right"),
    "l": ("l", "left"),
    "ctr": ("ctr", "center"),
    "r": ("r", "right"),
}
TEXT_VERTICAL_ALIGNMENTS = {
    "top": ("t", "top"),
    "middle": ("ctr", "middle"),
    "center": ("ctr", "middle"),
    "mid": ("ctr", "middle"),
    "centre": ("ctr", "middle"),
    "bottom": ("b", "bottom"),
    "t": ("t", "top"),
    "ctr": ("ctr", "middle"),
    "b": ("b", "bottom"),
}


def emu(value):
    return int(round(float(value) * EMU_PER_INCH))


def hex_color(value, default="000000"):
    if not value:
        return default
    return str(value).strip().lstrip("#").upper()


def content_type_for(path):
    suffix = Path(path).suffix.lower()
    if suffix == ".png":
        return "image/png"
    if suffix in (".jpg", ".jpeg"):
        return "image/jpeg"
    if suffix == ".gif":
        return "image/gif"
    if suffix == ".svg":
        return "image/svg+xml"
    raise ValueError(f"Unsupported image type: {path}")


def image_ext(path):
    suffix = Path(path).suffix.lower()
    return ".jpg" if suffix == ".jpeg" else suffix


def xml_text(value):
    return html.escape(str(value), quote=True)


def text_alignment(value="left"):
    key = str(value or "left").strip().lower()
    if key not in TEXT_ALIGNMENTS:
        raise ValueError(f"Unsupported text align value: {value}")
    return TEXT_ALIGNMENTS[key]


def text_vertical_alignment(value="top"):
    key = str(value or "top").strip().lower()
    if key not in TEXT_VERTICAL_ALIGNMENTS:
        raise ValueError(f"Unsupported text valign value: {value}")
    return TEXT_VERTICAL_ALIGNMENTS[key]


def source_size_px(manifest):
    source = manifest.get("source", {})
    width = source.get("width_px")
    height = source.get("height_px")
    if width and height:
        return float(width), float(height)
    return None


def slide_size(manifest):
    slide = manifest.get("slide", {})
    return float(slide.get("width", 13.333)), float(slide.get("height", 7.5))


def fit_content_box(source_width, source_height, slide_width, slide_height):
    source_aspect = source_width / source_height
    slide_aspect = slide_width / slide_height
    if source_aspect >= slide_aspect:
        width = slide_width
        height = width / source_aspect
        left = 0
        top = (slide_height - height) / 2
    else:
        height = slide_height
        width = height * source_aspect
        left = (slide_width - width) / 2
        top = 0
    return {"left": left, "top": top, "width": width, "height": height}


def content_box_for_manifest(manifest):
    content_box = manifest.get("content_box")
    if content_box:
        return {
            "left": float(content_box.get("left", 0)),
            "top": float(content_box.get("top", 0)),
            "width": float(content_box.get("width", 1)),
            "height": float(content_box.get("height", 1)),
        }
    source_size = source_size_px(manifest)
    slide_width, slide_height = slide_size(manifest)
    if not source_size:
        return {"left": 0, "top": 0, "width": slide_width, "height": slide_height}
    source_width, source_height = source_size
    return fit_content_box(source_width, source_height, slide_width, slide_height)


def px_to_inches(manifest, x, y, width, height):
    source_size = source_size_px(manifest)
    if not source_size:
        raise ValueError("Manifest uses pixel coordinates but lacks source.width_px/source.height_px")
    source_width, source_height = source_size
    content_box = content_box_for_manifest(manifest)
    return {
        "left": content_box["left"] + float(x) / source_width * content_box["width"],
        "top": content_box["top"] + float(y) / source_height * content_box["height"],
        "width": float(width) / source_width * content_box["width"],
        "height": float(height) / source_height * content_box["height"],
    }


def normalize_position_item(manifest, item):
    item = dict(item)
    if "polygon_px" in item:
        points = [(float(point[0]), float(point[1])) for point in item["polygon_px"]]
        if points and "box_px" not in item:
            xs = [point[0] for point in points]
            ys = [point[1] for point in points]
            item["box_px"] = [min(xs), min(ys), max(xs) - min(xs), max(ys) - min(ys)]
        item["polygon"] = [
            [px_to_inches(manifest, point[0], point[1], 0, 0)["left"], px_to_inches(manifest, point[0], point[1], 0, 0)["top"]]
            for point in points
        ]
    if "box_px" in item:
        x, y, width, height = item["box_px"]
        item.update(px_to_inches(manifest, x, y, width, height))
    if "points_px" in item:
        x1, y1, x2, y2 = item["points_px"]
        left = min(float(x1), float(x2))
        top = min(float(y1), float(y2))
        width = abs(float(x2) - float(x1))
        height = abs(float(y2) - float(y1))
        item.update(px_to_inches(manifest, left, top, width, height))
        start = px_to_inches(manifest, x1, y1, 0, 0)
        end = px_to_inches(manifest, x2, y2, 0, 0)
        item["points"] = [start["left"], start["top"], end["left"], end["top"]]
        if float(x2) < float(x1):
            item["flip_h"] = True
        if float(y2) < float(y1):
            item["flip_v"] = True
    if item.get("source_corner_radius_px") is not None and "radius" not in item:
        radius = float(item.get("source_corner_radius_px") or 0)
        item["radius"] = px_to_inches(manifest, 0, 0, radius, radius)["width"]
    return item


def normalize_font_size_item(manifest, item):
    """Resolve optional source-pixel font sizes to PowerPoint points."""

    item = dict(item)
    source = source_size_px(manifest)
    if source:
        _source_width, source_height = source
        content_box = content_box_for_manifest(manifest)

        def resolve(record):
            if not isinstance(record, dict) or record.get("font_size_px") is None:
                return
            if record.get("font_size") is not None:
                raise ValueError("provide font_size or font_size_px, not both")
            record["font_size"] = round(
                float(record.pop("font_size_px")) * float(content_box["height"]) * 72.0 / source_height,
                3,
            )

        resolve(item)
        for run in item.get("runs", []):
            resolve(run)
        for paragraph in item.get("paragraphs", []):
            if isinstance(paragraph, dict):
                resolve(paragraph)
                for run in paragraph.get("runs", []):
                    resolve(run)
        for row in item.get("rows", []):
            for cell in row:
                if not isinstance(cell, dict):
                    continue
                resolve(cell)
                for run in cell.get("runs", []):
                    resolve(run)
    return item


def iter_text_lines(item):
    if item.get("paragraphs"):
        lines = []
        for paragraph in item["paragraphs"]:
            if isinstance(paragraph, str):
                lines.append(paragraph)
            else:
                runs = paragraph.get("runs")
                if runs:
                    lines.append("".join(str(run.get("text", "")) for run in runs))
                else:
                    lines.append(str(paragraph.get("text", "")))
        return lines or [""]
    if item.get("runs"):
        return ["".join(str(run.get("text", "")) for run in item["runs"])]
    return str(item.get("text", "")).splitlines() or [""]


def text_width_units(text):
    units = 0.0
    for char in str(text):
        codepoint = ord(char)
        if char.isspace():
            units += 0.32
        elif codepoint <= 0x7F:
            units += 0.55
        elif 0xFF00 <= codepoint <= 0xFFEF:
            units += 1.0
        elif 0x4E00 <= codepoint <= 0x9FFF:
            units += 1.0
        else:
            units += 0.85
    return max(units, 1.0)


def longest_unbreakable_units(text):
    tokens = [token for token in re.split(r"\s+", str(text)) if token]
    if not tokens:
        return text_width_units(text)
    return max(text_width_units(token) for token in tokens)


def is_measured_text(item):
    """True when the author marked this box as sized from `page hints` measurement."""
    return str(item.get("font_size_source", "")).strip().lower() in {"measured", "hints"}


def resolve_text_fonts(item):
    """Bind text to an installed family before PowerPoint can substitute it."""

    if item.get("preserve_font_name"):
        return item

    def resolve(record, inherited=""):
        requested = str(record.get("font") or inherited or "PingFang SC").strip()
        path, family = find_font(requested)
        if family:
            if family.casefold() != requested.casefold():
                record.setdefault("_requested_font", requested)
            record["font"] = family
        if path:
            record["_resolved_font_path"] = str(path)

    resolve(item)
    base_font = str(item.get("font") or "PingFang SC")
    for run in item.get("runs", []):
        if isinstance(run, dict):
            resolve(run, base_font)
    for paragraph in item.get("paragraphs", []):
        if isinstance(paragraph, dict):
            for run in paragraph.get("runs", []):
                if isinstance(run, dict):
                    resolve(run, base_font)
    return item


def _font_at_em(path):
    if not path:
        return None
    try:
        return ImageFont.truetype(str(path), size=100)
    except OSError:
        return None


def _actual_text_limits(item, manifest):
    """Return width/height font-size limits using the resolved local font."""

    font = _font_at_em(item.get("_resolved_font_path"))
    if font is None or "width" not in item or "height" not in item:
        return None
    lines = iter_text_lines(item)
    width_pt = max(1.0, float(item["width"]) * 72)
    height_pt = max(1.0, float(item["height"]) * 72)
    requested = float(item.get("font_size", 18))
    widths_per_pt = [float(font.getlength(line)) / 100.0 for line in lines]
    wrap_enabled = item.get("wrap") not in (None, "", "none")
    if wrap_enabled:
        rendered_lines = sum(
            max(1, math.ceil((units * requested) / width_pt))
            for units in widths_per_pt
        )
        width_limit = requested
    else:
        rendered_lines = max(1, len(lines))
        width_limit = width_pt / max(max(widths_per_pt, default=1.0), 0.01)
    line_height = float(item.get("line_height", manifest.get("text_line_height", DEFAULT_TEXT_LINE_HEIGHT)))
    height_limit = height_pt / (rendered_lines * max(line_height, 1.0))
    return width_limit, height_limit


def fitted_font_size(item, manifest):
    if item.get("fit_text") is False or manifest.get("fit_text") is False:
        return None
    if "width" not in item or "height" not in item:
        return None
    actual_limits = _actual_text_limits(item, manifest)
    lines = iter_text_lines(item)
    requested = float(item.get("font_size", 18))
    width_pt = max(1.0, float(item.get("width", 1)) * 72)
    height_pt = max(1.0, float(item.get("height", 0.4)) * 72)
    if is_measured_text(item):
        # Box and font size were both measured from source ink; the safety
        # discount exists to absorb estimation error in hand-written boxes
        # and would systematically shrink correct text here. Clamp only at
        # the geometric limit.
        safety = 1.0
    else:
        safety = float(item.get("text_fit_safety", manifest.get("text_fit_safety", DEFAULT_TEXT_FIT_SAFETY)))
    line_height = float(item.get("line_height", manifest.get("text_line_height", DEFAULT_TEXT_LINE_HEIGHT)))
    wrap_enabled = item.get("wrap") not in (None, "", "none")
    if wrap_enabled:
        line_count = sum(max(1, math.ceil(text_width_units(line) * requested / width_pt)) for line in lines)
        width_limit = width_pt / max(longest_unbreakable_units(line) for line in lines)
    else:
        line_count = max(1, len(lines))
        width_limit = width_pt / max(text_width_units(line) for line in lines)
    height_limit = height_pt / (line_count * max(line_height, 1.0))
    max_font_size = min(width_limit, height_limit) * safety
    if actual_limits is not None:
        actual_safety = 0.98 if is_measured_text(item) else max(safety, 0.94)
        max_font_size = min(max_font_size, min(actual_limits) * actual_safety)
    explicit_max = item.get("max_font_size")
    if explicit_max not in (None, ""):
        max_font_size = min(max_font_size, float(explicit_max))
    min_font_size = float(item.get("min_font_size", manifest.get("min_font_size", DEFAULT_MIN_FONT_SIZE)))
    return max(min_font_size, max_font_size)


def scale_run_font_sizes(item, ratio):
    def scale_run(run):
        if run.get("font_size") not in (None, ""):
            run["font_size"] = round(float(run["font_size"]) * ratio, 1)

    for run in item.get("runs", []):
        scale_run(run)
    for paragraph in item.get("paragraphs", []):
        if isinstance(paragraph, dict):
            for run in paragraph.get("runs", []):
                scale_run(run)


def fit_text_item(item, manifest):
    fitted = fitted_font_size(item, manifest)
    if fitted is None:
        return item
    requested = float(item.get("font_size", fitted))
    effective = min(requested, fitted)
    if effective < requested:
        item["_requested_font_size"] = requested
        item["font_size"] = round(effective, 1)
        scale_run_font_sizes(item, effective / requested)
    elif "font_size" not in item:
        item["font_size"] = round(effective, 1)
    return item


def normalize_manifest(manifest):
    """Return a manifest copy with pixel authoring fields resolved to inches."""
    normalized = deepcopy(manifest)
    text_boxes = list(normalized.get("text_boxes", []))
    normalized["text_boxes"] = []
    for item in text_boxes:
        prepared = resolve_text_fonts(normalize_position_item(normalized, normalize_font_size_item(normalized, item)))
        normalized["text_boxes"].append(fit_text_item(prepared, normalized))
    for key in ("images", "shapes"):
        normalized[key] = [normalize_position_item(normalized, item) for item in normalized.get(key, [])]
    tables = list(normalized.get("tables", []))
    normalized["tables"] = []
    for item in tables:
        prepared = resolve_text_fonts(normalize_position_item(normalized, normalize_font_size_item(normalized, item)))
        for row in prepared.get("rows", []):
            for cell in row:
                if isinstance(cell, dict):
                    resolve_text_fonts(cell)
        normalized["tables"].append(prepared)
    return normalized


def preview_color(value):
    if not value or value == "none":
        return value
    value = str(value).strip()
    if value.startswith("#"):
        return value
    if len(value) == 6 and all(ch in "0123456789abcdefABCDEF" for ch in value):
        return f"#{value}"
    return value


def shape_fill(fill):
    if not fill or fill == "none":
        return '<a:noFill/>'
    return f'<a:solidFill><a:srgbClr val="{hex_color(fill)}"/></a:solidFill>'


def shape_fill_for_item(item):
    gradient = item.get("gradient")
    if not gradient:
        return shape_fill(item.get("fill"))
    stops = gradient.get("stops") if isinstance(gradient, dict) else None
    if not isinstance(stops, list) or len(stops) < 2:
        raise ValueError("gradient requires at least two color stops")
    stop_xml = []
    for stop in stops:
        position = float(stop.get("position", 0))
        if 0 <= position <= 1:
            position *= 100000
        if not 0 <= position <= 100000:
            raise ValueError("gradient stop position must be 0..1 or 0..100000")
        stop_xml.append(
            f'<a:gs pos="{int(round(position))}"><a:srgbClr val="{hex_color(stop.get("color"))}"/></a:gs>'
        )
    angle = int(round(float(gradient.get("angle", 0)) * 60000))
    return (
        '<a:gradFill rotWithShape="1"><a:gsLst>'
        + "".join(stop_xml)
        + f'</a:gsLst><a:lin ang="{angle}" scaled="1"/></a:gradFill>'
    )


def shape_line_xml(stroke, width, dash=None, *, start_arrow="none", end_arrow="none"):
    if not stroke or stroke == "none":
        return '<a:ln><a:noFill/></a:ln>'
    dash_xml = f'<a:prstDash val="{xml_text(dash)}"/>' if dash else ""
    head_xml = f'<a:headEnd type="{xml_text(start_arrow)}"/>' if start_arrow and start_arrow != "none" else ""
    tail_xml = f'<a:tailEnd type="{xml_text(end_arrow)}"/>' if end_arrow and end_arrow != "none" else ""
    return (
        f'<a:ln w="{int(float(width or 1) * 12700)}">'
        f'<a:solidFill><a:srgbClr val="{hex_color(stroke)}"/></a:solidFill>'
        f"{dash_xml}{head_xml}{tail_xml}"
        "</a:ln>"
    )


def slide_background_xml(slide):
    background = slide.get("background")
    if not background:
        return ""
    return (
        "<p:bg><p:bgPr>"
        f'<a:solidFill><a:srgbClr val="{hex_color(background, "FFFFFF")}"/></a:solidFill>'
        "<a:effectLst/></p:bgPr></p:bg>"
    )


def text_box_xml(idx, item):
    left = emu(item.get("left", 0))
    top = emu(item.get("top", 0))
    width = emu(item.get("width", 1))
    height = emu(item.get("height", 0.4))
    rotation = item.get("rotation")
    rotation_attr = f' rot="{int(float(rotation) * 60000)}"' if rotation not in (None, "") else ""
    font_size = int(float(item.get("font_size", 18)) * 100)
    font = xml_text(item.get("font", "PingFang SC"))
    align = text_alignment(item.get("align", "left"))[0]
    anchor = text_vertical_alignment(item.get("valign", "top"))[0]
    wrap = item.get("wrap", "none")
    autofit = item.get("autofit", "none")
    autofit_xml = "<a:spAutoFit/>" if autofit == "shape" else "<a:noAutofit/>"
    paragraphs = item.get("paragraphs")
    runs = item.get("runs")

    def run_xml(run):
        run_font_size = int(float(run.get("font_size", item.get("font_size", 18))) * 100)
        run_font = xml_text(run.get("font", item.get("font", "PingFang SC")))
        run_color = hex_color(run.get("color", item.get("color", "#111111")))
        run_bold = ' b="1"' if run.get("bold", item.get("bold")) else ""
        run_italic = ' i="1"' if run.get("italic", item.get("italic")) else ""
        run_baseline = run.get("baseline")
        run_baseline_attr = f' baseline="{int(float(run_baseline))}"' if run_baseline not in (None, "") else ""
        run_text = xml_text(run.get("text", ""))
        return (
            f'<a:r><a:rPr lang="zh-CN" sz="{run_font_size}"{run_bold}{run_italic}{run_baseline_attr}>'
            f'<a:solidFill><a:srgbClr val="{run_color}"/></a:solidFill>'
            f'<a:latin typeface="{run_font}"/><a:ea typeface="{run_font}"/><a:cs typeface="{run_font}"/>'
            f'</a:rPr><a:t>{run_text}</a:t></a:r>'
        )

    def paragraph_xml(paragraph):
        if isinstance(paragraph, str):
            paragraph_runs = [{"text": paragraph}]
        else:
            paragraph_runs = paragraph.get("runs", [{"text": paragraph.get("text", "")}])
        return (
            f'<a:p><a:pPr algn="{align}"/>'
            + "".join(run_xml(run) for run in paragraph_runs)
            + f'<a:endParaRPr lang="zh-CN" sz="{font_size}"/></a:p>'
        )

    if paragraphs:
        text_body = "".join(paragraph_xml(paragraph) for paragraph in paragraphs)
    elif runs:
        text_body = paragraph_xml({"runs": runs})
    else:
        text_body = "".join(paragraph_xml(part) for part in str(item.get("text", "")).split("\n"))
    return f"""
      <p:sp>
        <p:nvSpPr><p:cNvPr id="{idx}" name="TextBox {idx}"/><p:cNvSpPr txBox="1"/><p:nvPr/></p:nvSpPr>
        <p:spPr><a:xfrm{rotation_attr}><a:off x="{left}" y="{top}"/><a:ext cx="{width}" cy="{height}"/></a:xfrm><a:prstGeom prst="rect"><a:avLst/></a:prstGeom><a:noFill/><a:ln><a:noFill/></a:ln></p:spPr>
        <p:txBody>
          <a:bodyPr wrap="{xml_text(wrap)}" anchor="{anchor}" lIns="0" tIns="0" rIns="0" bIns="0">{autofit_xml}</a:bodyPr><a:lstStyle/>
          {text_body}
        </p:txBody>
      </p:sp>"""


def image_xml(idx, rel_id, item):
    left = emu(item.get("left", 0))
    top = emu(item.get("top", 0))
    width = emu(item.get("width", 1))
    height = emu(item.get("height", 1))
    name = xml_text(item.get("alt") or Path(item.get("path", "")).stem or f"Image {idx}")
    return f"""
      <p:pic>
        <p:nvPicPr><p:cNvPr id="{idx}" name="{name}"/><p:cNvPicPr/><p:nvPr/></p:nvPicPr>
        <p:blipFill><a:blip r:embed="{rel_id}"/><a:stretch><a:fillRect/></a:stretch></p:blipFill>
        <p:spPr><a:xfrm><a:off x="{left}" y="{top}"/><a:ext cx="{width}" cy="{height}"/></a:xfrm><a:prstGeom prst="rect"><a:avLst/></a:prstGeom></p:spPr>
      </p:pic>"""


def _table_text_xml(value, table, row_index):
    cell = value if isinstance(value, dict) else {"text": value}
    font_size = int(float(cell.get("font_size", table.get("font_size", 12))) * 100)
    bold = cell.get("bold", table.get("header_bold", row_index == 0))
    align = text_alignment(cell.get("align", table.get("align", "center")))[0]
    valign = text_vertical_alignment(cell.get("valign", table.get("valign", "middle")))[0]
    raw_runs = cell.get("runs")
    if not raw_runs:
        raw_runs = [{"text": cell.get("text", "")}]
    runs = []
    for run in raw_runs:
        run = run if isinstance(run, dict) else {"text": run}
        run_size = int(float(run.get("font_size", cell.get("font_size", table.get("font_size", 12)))) * 100)
        run_font = xml_text(run.get("font", cell.get("font", table.get("font", "PingFang SC"))))
        run_color = hex_color(run.get("color", cell.get("color", table.get("header_color" if row_index == 0 else "color", "#111111"))))
        run_bold = run.get("bold", bold)
        run_bold_attr = ' b="1"' if run_bold else ""
        runs.append(
            f'<a:r><a:rPr lang="zh-CN" sz="{run_size}"{run_bold_attr}>'
            f'<a:solidFill><a:srgbClr val="{run_color}"/></a:solidFill>'
            f'<a:latin typeface="{run_font}"/><a:ea typeface="{run_font}"/><a:cs typeface="{run_font}"/>'
            f'</a:rPr><a:t>{xml_text(run.get("text", ""))}</a:t></a:r>'
        )
    return (
        f'<a:txBody><a:bodyPr wrap="square" anchor="{valign}" lIns="45720" tIns="22860" rIns="45720" bIns="22860"/>'
        f'<a:lstStyle/><a:p><a:pPr algn="{align}"/>{"".join(runs)}'
        f'<a:endParaRPr lang="zh-CN" sz="{font_size}"/></a:p></a:txBody>'
    )


def _table_cell_xml(value, table, row_index):
    cell = value if isinstance(value, dict) else {}
    fill = cell.get("fill", table.get("header_fill" if row_index == 0 else "fill", "#FFFFFF"))
    border = cell.get("border", table.get("border", "#B7C4D4"))
    border_width = int(float(cell.get("border_width", table.get("border_width", 0.75))) * 12700)
    border_xml = "".join(
        f'<a:ln{side} w="{border_width}"><a:solidFill><a:srgbClr val="{hex_color(border)}"/></a:solidFill></a:ln{side}>'
        for side in ("L", "R", "T", "B")
    ) if border and border != "none" else "".join(f'<a:ln{side}><a:noFill/></a:ln{side}>' for side in ("L", "R", "T", "B"))
    fill_xml = shape_fill(fill)
    return f'<a:tc>{_table_text_xml(value, table, row_index)}<a:tcPr>{border_xml}{fill_xml}</a:tcPr></a:tc>'


def table_xml(idx, item):
    rows = item.get("rows") or []
    if not rows or not isinstance(rows, list):
        raise ValueError("table rows must be a non-empty list")
    column_count = max(len(row) for row in rows)
    if column_count <= 0:
        raise ValueError("table requires at least one column")
    left = emu(item.get("left", 0))
    top = emu(item.get("top", 0))
    width = emu(item.get("width", 1))
    height = emu(item.get("height", 1))
    column_weights = [float(value) for value in item.get("column_widths", [1] * column_count)]
    if len(column_weights) != column_count or sum(column_weights) <= 0:
        raise ValueError("column_widths must match the table column count")
    grid = [int(width * value / sum(column_weights)) for value in column_weights]
    grid[-1] += width - sum(grid)
    row_weights = [float(value) for value in item.get("row_heights", [1] * len(rows))]
    if len(row_weights) != len(rows) or sum(row_weights) <= 0:
        raise ValueError("row_heights must match the table row count")
    row_sizes = [int(height * value / sum(row_weights)) for value in row_weights]
    row_sizes[-1] += height - sum(row_sizes)
    grid_xml = "".join(f'<a:gridCol w="{value}"/>' for value in grid)
    rows_xml = []
    for row_index, (row, row_height) in enumerate(zip(rows, row_sizes)):
        values = list(row) + [""] * (column_count - len(row))
        rows_xml.append(
            f'<a:tr h="{row_height}">' + "".join(_table_cell_xml(value, item, row_index) for value in values) + '</a:tr>'
        )
    name = xml_text(item.get("name", f"Table {idx}"))
    return f"""
      <p:graphicFrame>
        <p:nvGraphicFramePr><p:cNvPr id="{idx}" name="{name}"/><p:cNvGraphicFramePr/><p:nvPr/></p:nvGraphicFramePr>
        <p:xfrm><a:off x="{left}" y="{top}"/><a:ext cx="{width}" cy="{height}"/></p:xfrm>
        <a:graphic><a:graphicData uri="http://schemas.openxmlformats.org/drawingml/2006/table">
          <a:tbl><a:tblPr firstRow="1" bandRow="0"/><a:tblGrid>{grid_xml}</a:tblGrid>{''.join(rows_xml)}</a:tbl>
        </a:graphicData></a:graphic>
      </p:graphicFrame>"""


def shape_xml(idx, item):
    kind = item.get("type", "rect")
    left = emu(item.get("left", 0))
    top = emu(item.get("top", 0))
    width = emu(item.get("width", 1))
    height = emu(item.get("height", 1))
    stroke_width = item.get("stroke_width", 1)
    flip_h = ' flipH="1"' if item.get("flip_h") else ""
    flip_v = ' flipV="1"' if item.get("flip_v") else ""
    fill = shape_fill_for_item(item)
    line = shape_line_xml(
        item.get("stroke", "#000000"),
        stroke_width,
        item.get("dash"),
        start_arrow=item.get("start_arrow", item.get("begin_arrow", "none")),
        end_arrow=item.get("end_arrow", "none"),
    )
    preset = item.get("preset")
    if item.get("polygon_px"):
        geometry = custom_polygon_geometry_xml(item)
    else:
        if not preset:
            preset = "line" if kind == "line" else "ellipse" if kind == "ellipse" else "roundRect" if kind == "roundRect" else "rect"
        geometry = preset_geometry_xml(preset, item)
    if kind == "line" or str(preset or "").endswith("Connector3"):
        return f"""
      <p:cxnSp>
        <p:nvCxnSpPr><p:cNvPr id="{idx}" name="Connector {idx}"/><p:cNvCxnSpPr/><p:nvPr/></p:nvCxnSpPr>
        <p:spPr><a:xfrm{flip_h}{flip_v}><a:off x="{left}" y="{top}"/><a:ext cx="{width}" cy="{height}"/></a:xfrm>{geometry}{fill}{line}</p:spPr>
      </p:cxnSp>"""
    return f"""
      <p:sp>
        <p:nvSpPr><p:cNvPr id="{idx}" name="{xml_text(kind.title())} {idx}"/><p:cNvSpPr/><p:nvPr/></p:nvSpPr>
        <p:spPr><a:xfrm{flip_h}{flip_v}><a:off x="{left}" y="{top}"/><a:ext cx="{width}" cy="{height}"/></a:xfrm>{geometry}{fill}{line}</p:spPr>
        <p:txBody><a:bodyPr/><a:lstStyle/><a:p/></p:txBody>
      </p:sp>"""


def custom_polygon_geometry_xml(item):
    points = [(float(point[0]), float(point[1])) for point in item.get("polygon_px", [])]
    if len(points) < 3:
        return '<a:prstGeom prst="rect"><a:avLst/></a:prstGeom>'
    box = item.get("box_px")
    if box and len(box) == 4:
        left, top, width, height = [float(value) for value in box]
    else:
        xs = [point[0] for point in points]
        ys = [point[1] for point in points]
        left, top = min(xs), min(ys)
        width, height = max(xs) - left, max(ys) - top
    width = max(width, 1.0)
    height = max(height, 1.0)

    def rel_coord(point):
        x, y = point
        return int(round((x - left) / width * 21600)), int(round((y - top) / height * 21600))

    first_x, first_y = rel_coord(points[0])
    segments = [f'<a:moveTo><a:pt x="{first_x}" y="{first_y}"/></a:moveTo>']
    for point in points[1:]:
        x, y = rel_coord(point)
        segments.append(f'<a:lnTo><a:pt x="{x}" y="{y}"/></a:lnTo>')
    segments.append("<a:close/>")
    return (
        '<a:custGeom><a:avLst/><a:gdLst/><a:ahLst/><a:cxnLst/>'
        '<a:rect l="l" t="t" r="r" b="b"/>'
        '<a:pathLst><a:path w="21600" h="21600">'
        + "".join(segments)
        + "</a:path></a:pathLst></a:custGeom>"
    )


def round_rect_adjustment(item):
    box_px = item.get("box_px")
    if item.get("source_corner_radius_px") is not None and box_px and len(box_px) == 4:
        radius = float(item.get("source_corner_radius_px") or 0)
        min_dim = max(1.0, min(float(box_px[2]), float(box_px[3])))
    elif item.get("radius") is not None and box_px and len(box_px) == 4:
        # A manifest authored in source pixels naturally expresses its corner
        # radius in the same coordinate space.  Treating this value as inches
        # turns a small radius into a capsule or arc on large containers.
        radius = float(item.get("radius") or 0)
        min_dim = max(1.0, min(float(box_px[2]), float(box_px[3])))
    elif item.get("radius") is not None:
        radius = float(item.get("radius") or 0)
        min_dim = max(0.01, min(float(item.get("width", 1)), float(item.get("height", 1))))
    else:
        return None
    return max(0, min(50000, int(round(radius / min_dim * 100000))))


def preset_geometry_xml(preset, item):
    if preset != "roundRect":
        return f'<a:prstGeom prst="{preset}"><a:avLst/></a:prstGeom>'
    adjustment = round_rect_adjustment(item)
    if adjustment is None:
        return '<a:prstGeom prst="roundRect"><a:avLst/></a:prstGeom>'
    return (
        '<a:prstGeom prst="roundRect"><a:avLst>'
        f'<a:gd name="adj" fmla="val {adjustment}"/>'
        '</a:avLst></a:prstGeom>'
    )


def slide_xml(manifest):
    slide = manifest.get("slide", {})
    next_id = 2
    parts = []
    layered = []
    for index, item in enumerate(manifest.get("shapes", [])):
        layered.append((float(item.get("z_index", 100)), index, "shape", item, None))
    for rel_index, item in enumerate(manifest.get("images", []), start=1):
        layered.append((float(item.get("z_index", 200)), rel_index, "image", item, f"rId{rel_index + 1}"))
    for index, item in enumerate(manifest.get("tables", [])):
        layered.append((float(item.get("z_index", 250)), index, "table", item, None))
    for index, item in enumerate(manifest.get("text_boxes", [])):
        layered.append((float(item.get("z_index", 300)), index, "text", item, None))

    for _z_index, _order, kind, item, rel_id in sorted(layered, key=lambda entry: (entry[0], entry[1])):
        if kind == "shape":
            parts.append(shape_xml(next_id, item))
        elif kind == "image":
            parts.append(image_xml(next_id, rel_id, item))
        elif kind == "table":
            parts.append(table_xml(next_id, item))
        else:
            parts.append(text_box_xml(next_id, item))
        next_id += 1
    return f"""<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<p:sld xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships" xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">
  <p:cSld>
    {slide_background_xml(slide)}
    <p:spTree>
      <p:nvGrpSpPr><p:cNvPr id="1" name=""/><p:cNvGrpSpPr/><p:nvPr/></p:nvGrpSpPr>
      <p:grpSpPr><a:xfrm><a:off x="0" y="0"/><a:ext cx="0" cy="0"/><a:chOff x="0" y="0"/><a:chExt cx="0" cy="0"/></a:xfrm></p:grpSpPr>
      {''.join(parts)}
    </p:spTree>
  </p:cSld>
  <p:clrMapOvr><a:masterClrMapping/></p:clrMapOvr>
</p:sld>"""


def rels_xml(manifest, media_start=1, notes_index=None):
    rels = ['<Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideLayout" Target="../slideLayouts/slideLayout1.xml"/>']
    for i, item in enumerate(manifest.get("images", []), start=1):
        target = f"../media/image{media_start + i - 1}{image_ext(item['path'])}"
        rels.append(f'<Relationship Id="rId{i + 1}" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/image" Target="{target}"/>')
    if notes_index is not None:
        rels.append(
            f'<Relationship Id="rId{len(rels) + 1}" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/notesSlide" Target="../notesSlides/notesSlide{notes_index}.xml"/>'
        )
    return '<?xml version="1.0" encoding="UTF-8" standalone="yes"?><Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">' + "".join(rels) + "</Relationships>"


def notes_slide_xml(text):
    paras = "".join(
        f'<a:p><a:r><a:rPr lang="zh-CN" sz="1200"/><a:t>{xml_text(line)}</a:t></a:r><a:endParaRPr lang="zh-CN" sz="1200"/></a:p>'
        for line in str(text).splitlines()
    ) or '<a:p><a:endParaRPr lang="zh-CN" sz="1200"/></a:p>'
    return f"""<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<p:notes xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships" xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">
  <p:cSld>
    <p:spTree>
      <p:nvGrpSpPr><p:cNvPr id="1" name=""/><p:cNvGrpSpPr/><p:nvPr/></p:nvGrpSpPr>
      <p:grpSpPr><a:xfrm><a:off x="0" y="0"/><a:ext cx="0" cy="0"/><a:chOff x="0" y="0"/><a:chExt cx="0" cy="0"/></a:xfrm></p:grpSpPr>
      <p:sp>
        <p:nvSpPr><p:cNvPr id="2" name="Notes Placeholder"/><p:cNvSpPr txBox="1"/><p:nvPr><p:ph type="body" idx="1"/></p:nvPr></p:nvSpPr>
        <p:spPr><a:xfrm><a:off x="685800" y="914400"/><a:ext cx="5486400" cy="6858000"/></a:xfrm><a:prstGeom prst="rect"><a:avLst/></a:prstGeom><a:noFill/><a:ln><a:noFill/></a:ln></p:spPr>
        <p:txBody><a:bodyPr/><a:lstStyle/>{paras}</p:txBody>
      </p:sp>
    </p:spTree>
  </p:cSld>
  <p:clrMapOvr><a:masterClrMapping/></p:clrMapOvr>
</p:notes>"""


def notes_rels_xml(slide_index):
    return f"""<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Relationships xmlns="{REL_NS}">
  <Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/notesMaster" Target="../notesMasters/notesMaster1.xml"/>
  <Relationship Id="rId2" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide" Target="../slides/slide{slide_index}.xml"/>
</Relationships>"""


def notes_master_xml():
    return """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<p:notesMaster xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships" xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">
  <p:cSld><p:spTree><p:nvGrpSpPr><p:cNvPr id="1" name=""/><p:cNvGrpSpPr/><p:nvPr/></p:nvGrpSpPr><p:grpSpPr/></p:spTree></p:cSld>
  <p:clrMap bg1="lt1" tx1="dk1" bg2="lt2" tx2="dk2" accent1="accent1" accent2="accent2" accent3="accent3" accent4="accent4" accent5="accent5" accent6="accent6" hlink="hlink" folHlink="folHlink"/>
</p:notesMaster>"""


def notes_master_rels_xml():
    return f"""<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Relationships xmlns="{REL_NS}">
  <Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/theme" Target="../theme/theme1.xml"/>
</Relationships>"""


def content_types_xml(manifests, notes_indices=None):
    notes_indices = notes_indices or []
    defaults = {
        "rels": "application/vnd.openxmlformats-package.relationships+xml",
        "xml": "application/xml",
    }
    for manifest in manifests:
        for item in manifest.get("images", []):
            ext = image_ext(item["path"]).lstrip(".")
            defaults[ext] = content_type_for(item["path"])
    default_xml = "".join(f'<Default Extension="{ext}" ContentType="{ctype}"/>' for ext, ctype in defaults.items())
    overrides = [
        '<Override PartName="/ppt/presentation.xml" ContentType="application/vnd.openxmlformats-officedocument.presentationml.presentation.main+xml"/>',
        '<Override PartName="/ppt/slideMasters/slideMaster1.xml" ContentType="application/vnd.openxmlformats-officedocument.presentationml.slideMaster+xml"/>',
        '<Override PartName="/ppt/slideLayouts/slideLayout1.xml" ContentType="application/vnd.openxmlformats-officedocument.presentationml.slideLayout+xml"/>',
        '<Override PartName="/ppt/theme/theme1.xml" ContentType="application/vnd.openxmlformats-officedocument.theme+xml"/>',
        '<Override PartName="/docProps/core.xml" ContentType="application/vnd.openxmlformats-package.core-properties+xml"/>',
        '<Override PartName="/docProps/app.xml" ContentType="application/vnd.openxmlformats-officedocument.extended-properties+xml"/>',
    ]
    for i in range(1, len(manifests) + 1):
        overrides.append(f'<Override PartName="/ppt/slides/slide{i}.xml" ContentType="application/vnd.openxmlformats-officedocument.presentationml.slide+xml"/>')
    if notes_indices:
        overrides.append('<Override PartName="/ppt/notesMasters/notesMaster1.xml" ContentType="application/vnd.openxmlformats-officedocument.presentationml.notesMaster+xml"/>')
        for i in notes_indices:
            overrides.append(f'<Override PartName="/ppt/notesSlides/notesSlide{i}.xml" ContentType="application/vnd.openxmlformats-officedocument.presentationml.notesSlide+xml"/>')
    return f'<?xml version="1.0" encoding="UTF-8" standalone="yes"?><Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">{default_xml}{"".join(overrides)}</Types>'


def is_wide_slide(width, height):
    return abs((float(width) / float(height)) / ASPECT_16_9 - 1) <= ASPECT_TOLERANCE


def slide_size_type(width, height):
    return "wide" if is_wide_slide(width, height) else "custom"


def presentation_xml(slide_count, width, height):
    slide_ids = "".join(f'<p:sldId id="{255 + i}" r:id="rId{i + 1}"/>' for i in range(1, slide_count + 1))
    size_type = slide_size_type(width, height)
    return f"""<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<p:presentation xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships" xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">
  <p:sldMasterIdLst><p:sldMasterId id="2147483648" r:id="rId1"/></p:sldMasterIdLst>
  <p:sldIdLst>{slide_ids}</p:sldIdLst>
  <p:sldSz cx="{width}" cy="{height}" type="{size_type}"/>
  <p:notesSz cx="6858000" cy="9144000"/>
</p:presentation>"""


def presentation_rels_xml(slide_count):
    rels = ['<Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideMaster" Target="slideMasters/slideMaster1.xml"/>']
    for i in range(1, slide_count + 1):
        rels.append(f'<Relationship Id="rId{i + 1}" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide" Target="slides/slide{i}.xml"/>')
    return f'<?xml version="1.0" encoding="UTF-8" standalone="yes"?><Relationships xmlns="{REL_NS}">{"".join(rels)}</Relationships>'


def write_common_parts(z, slide_count, width, height, notes_count):
    presentation_format = "Widescreen" if slide_size_type(width, height) == "wide" else "Custom"
    z.writestr("_rels/.rels", """<?xml version="1.0" encoding="UTF-8" standalone="yes"?><Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships"><Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" Target="ppt/presentation.xml"/><Relationship Id="rId2" Type="http://schemas.openxmlformats.org/package/2006/relationships/metadata/core-properties" Target="docProps/core.xml"/><Relationship Id="rId3" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/extended-properties" Target="docProps/app.xml"/></Relationships>""")
    z.writestr("docProps/core.xml", """<?xml version="1.0" encoding="UTF-8" standalone="yes"?><cp:coreProperties xmlns:cp="http://schemas.openxmlformats.org/package/2006/metadata/core-properties" xmlns:dc="http://purl.org/dc/elements/1.1/"><dc:title>Image to editable PPT</dc:title></cp:coreProperties>""")
    z.writestr("docProps/app.xml", f"""<?xml version="1.0" encoding="UTF-8" standalone="yes"?><Properties xmlns="http://schemas.openxmlformats.org/officeDocument/2006/extended-properties"><Application>Codex</Application><PresentationFormat>{presentation_format}</PresentationFormat><Slides>{slide_count}</Slides></Properties>""")
    z.writestr("ppt/presentation.xml", presentation_xml(slide_count, width, height))
    z.writestr("ppt/_rels/presentation.xml.rels", presentation_rels_xml(slide_count))
    z.writestr("ppt/slideMasters/slideMaster1.xml", """<?xml version="1.0" encoding="UTF-8" standalone="yes"?><p:sldMaster xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships" xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"><p:cSld><p:spTree><p:nvGrpSpPr><p:cNvPr id="1" name=""/><p:cNvGrpSpPr/><p:nvPr/></p:nvGrpSpPr><p:grpSpPr/></p:spTree></p:cSld><p:clrMap bg1="lt1" tx1="dk1" bg2="lt2" tx2="dk2" accent1="accent1" accent2="accent2" accent3="accent3" accent4="accent4" accent5="accent5" accent6="accent6" hlink="hlink" folHlink="folHlink"/><p:sldLayoutIdLst><p:sldLayoutId id="2147483649" r:id="rId1"/></p:sldLayoutIdLst></p:sldMaster>""")
    z.writestr("ppt/slideMasters/_rels/slideMaster1.xml.rels", """<?xml version="1.0" encoding="UTF-8" standalone="yes"?><Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships"><Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideLayout" Target="../slideLayouts/slideLayout1.xml"/><Relationship Id="rId2" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/theme" Target="../theme/theme1.xml"/></Relationships>""")
    z.writestr("ppt/slideLayouts/slideLayout1.xml", """<?xml version="1.0" encoding="UTF-8" standalone="yes"?><p:sldLayout xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships" xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" type="blank" preserve="1"><p:cSld name="Blank"><p:spTree><p:nvGrpSpPr><p:cNvPr id="1" name=""/><p:cNvGrpSpPr/><p:nvPr/></p:nvGrpSpPr><p:grpSpPr/></p:spTree></p:cSld><p:clrMapOvr><a:masterClrMapping/></p:clrMapOvr></p:sldLayout>""")
    z.writestr("ppt/slideLayouts/_rels/slideLayout1.xml.rels", """<?xml version="1.0" encoding="UTF-8" standalone="yes"?><Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships"><Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideMaster" Target="../slideMasters/slideMaster1.xml"/></Relationships>""")
    z.writestr("ppt/theme/theme1.xml", """<?xml version="1.0" encoding="UTF-8" standalone="yes"?><a:theme xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" name="ImageToEditablePPT"><a:themeElements><a:clrScheme name="Office"><a:dk1><a:sysClr val="windowText" lastClr="000000"/></a:dk1><a:lt1><a:sysClr val="window" lastClr="FFFFFF"/></a:lt1><a:dk2><a:srgbClr val="1F1F1F"/></a:dk2><a:lt2><a:srgbClr val="F8F8F8"/></a:lt2><a:accent1><a:srgbClr val="0F766E"/></a:accent1><a:accent2><a:srgbClr val="E66B00"/></a:accent2><a:accent3><a:srgbClr val="F6D365"/></a:accent3><a:accent4><a:srgbClr val="57C4B8"/></a:accent4><a:accent5><a:srgbClr val="666666"/></a:accent5><a:accent6><a:srgbClr val="111111"/></a:accent6><a:hlink><a:srgbClr val="0563C1"/></a:hlink><a:folHlink><a:srgbClr val="954F72"/></a:folHlink></a:clrScheme><a:fontScheme name="PingFang"><a:majorFont><a:latin typeface="PingFang SC"/><a:ea typeface="PingFang SC"/><a:cs typeface="PingFang SC"/></a:majorFont><a:minorFont><a:latin typeface="PingFang SC"/><a:ea typeface="PingFang SC"/><a:cs typeface="PingFang SC"/></a:minorFont></a:fontScheme><a:fmtScheme name="Office"><a:fillStyleLst><a:solidFill><a:schemeClr val="phClr"/></a:solidFill></a:fillStyleLst><a:lnStyleLst><a:ln w="9525"><a:solidFill><a:schemeClr val="phClr"/></a:solidFill></a:ln></a:lnStyleLst><a:effectStyleLst><a:effectStyle><a:effectLst/></a:effectStyle></a:effectStyleLst><a:bgFillStyleLst><a:solidFill><a:schemeClr val="phClr"/></a:solidFill></a:bgFillStyleLst></a:fmtScheme></a:themeElements></a:theme>""")
    if notes_count:
        z.writestr("ppt/notesMasters/notesMaster1.xml", notes_master_xml())
        z.writestr("ppt/notesMasters/_rels/notesMaster1.xml.rels", notes_master_rels_xml())


def normalize_powerpoint_package(out_path, slide_count, width, height, notes_count=0):
    """Wrap generated slide XML in a standard python-pptx package.

    The editable slide XML is intentionally authored directly because it keeps
    this Builder small and deterministic.  PowerPoint is stricter than
    python-pptx and LibreOffice about the surrounding presentation package,
    however, so use python-pptx only as the standards-compatible container and
    transplant our editable slides and media into it.
    """
    from pptx import Presentation

    out = Path(out_path)
    with zipfile.ZipFile(out) as package:
        generated = {name: package.read(name) for name in package.namelist()}

    with tempfile.TemporaryDirectory(prefix="editppt-ooxml-") as temporary:
        shell_path = Path(temporary) / "shell.pptx"
        presentation = Presentation()
        presentation.slide_width = int(width)
        presentation.slide_height = int(height)
        blank = presentation.slide_layouts[6]
        for _ in range(slide_count):
            presentation.slides.add_slide(blank)
        if notes_count:
            # Stable one-to-one notesSlide numbering lets source note parts be
            # transplanted without rewriting their slide references.
            for slide in presentation.slides:
                slide.notes_slide
        presentation.save(shell_path)
        with zipfile.ZipFile(shell_path) as package:
            standard = {name: package.read(name) for name in package.namelist()}

    relationship_tag = f"{{{REL_NS}}}Relationship"
    for slide_index in range(1, slide_count + 1):
        slide_name = f"ppt/slides/slide{slide_index}.xml"
        rels_name = f"ppt/slides/_rels/slide{slide_index}.xml.rels"
        standard[slide_name] = generated[slide_name]
        target_rels = ET.fromstring(standard[rels_name])
        source_rels = ET.fromstring(generated[rels_name])
        for source_rel in source_rels.findall(relationship_tag):
            if str(source_rel.get("Type") or "").endswith("/slideLayout"):
                continue
            for existing in list(target_rels):
                if existing.get("Id") == source_rel.get("Id"):
                    target_rels.remove(existing)
            target_rels.append(deepcopy(source_rel))
        standard[rels_name] = ET.tostring(
            target_rels,
            encoding="UTF-8",
            xml_declaration=True,
        )

    for name, data in generated.items():
        if name.startswith("ppt/media/") or name.startswith("ppt/notesSlides/"):
            standard[name] = data

    content_types_ns = "http://schemas.openxmlformats.org/package/2006/content-types"
    target_types = ET.fromstring(standard["[Content_Types].xml"])
    source_types = ET.fromstring(generated["[Content_Types].xml"])
    existing_types = {
        (child.tag, child.get("Extension") or child.get("PartName"))
        for child in target_types
    }
    for child in source_types:
        key = (child.tag, child.get("Extension") or child.get("PartName"))
        if key not in existing_types:
            target_types.append(deepcopy(child))
            existing_types.add(key)
    ET.register_namespace("", content_types_ns)
    standard["[Content_Types].xml"] = ET.tostring(
        target_types,
        encoding="UTF-8",
        xml_declaration=True,
    )

    app_ns = "http://schemas.openxmlformats.org/officeDocument/2006/extended-properties"
    app_properties = ET.fromstring(standard["docProps/app.xml"])
    slides_property = app_properties.find(f"{{{app_ns}}}Slides")
    if slides_property is not None:
        slides_property.text = str(slide_count)
    format_property = app_properties.find(f"{{{app_ns}}}PresentationFormat")
    if format_property is not None:
        format_property.text = "Widescreen" if slide_size_type(width, height) == "wide" else "Custom"
    ET.register_namespace("", app_ns)
    standard["docProps/app.xml"] = ET.tostring(
        app_properties,
        encoding="UTF-8",
        xml_declaration=True,
    )

    temporary_out = out.with_suffix(out.suffix + ".tmp")
    with zipfile.ZipFile(temporary_out, "w", zipfile.ZIP_DEFLATED) as package:
        for name, data in standard.items():
            package.writestr(name, data)
    temporary_out.replace(out)


def write_pptx(manifest, out_path, manifest_path):
    width = emu(manifest.get("slide", {}).get("width", 13.333))
    height = emu(manifest.get("slide", {}).get("height", 7.5))
    out = Path(out_path)
    out.parent.mkdir(parents=True, exist_ok=True)
    normalized = normalize_manifest(manifest)
    media_index = 1
    with zipfile.ZipFile(out, "w", zipfile.ZIP_DEFLATED) as z:
        z.writestr("[Content_Types].xml", content_types_xml([normalized], []))
        write_common_parts(z, 1, width, height, 0)
        z.writestr("ppt/slides/slide1.xml", slide_xml(normalized))
        z.writestr("ppt/slides/_rels/slide1.xml.rels", rels_xml(normalized, media_index, None))
        base = Path(manifest_path).resolve().parent
        for item in normalized.get("images", []):
            src = Path(item["path"])
            if not src.is_absolute():
                src = base / src
            z.write(src, f"ppt/media/image{media_index}{image_ext(src)}")
            media_index += 1
    normalize_powerpoint_package(out, 1, width, height)
    return normalized


def build_report(normalized, out_path):
    substitutions = []
    text_adjustments = []
    for index, item in enumerate(normalized.get("text_boxes", []), start=1):
        if item.get("_requested_font"):
            substitutions.append({
                "text_box": index,
                "requested": item["_requested_font"],
                "resolved": item.get("font", ""),
            })
        if item.get("_requested_font_size") is not None:
            requested = float(item["_requested_font_size"])
            effective = float(item.get("font_size", requested))
            text_adjustments.append({
                "text_box": index,
                "requested_pt": requested,
                "effective_pt": effective,
                "shrink_ratio": round(effective / max(requested, 0.01), 4),
            })
    severe_adjustments = [
        value for value in text_adjustments if float(value.get("shrink_ratio", 1)) < 0.75
    ]
    return {
        "schema_version": 1,
        "output_pptx": str(Path(out_path).resolve()),
        "objects": {
            "shapes": len(normalized.get("shapes", [])),
            "images": len(normalized.get("images", [])),
            "tables": len(normalized.get("tables", [])),
            "text_boxes": len(normalized.get("text_boxes", [])),
        },
        "font_substitutions": substitutions,
        "text_adjustments": text_adjustments,
        "severe_text_adjustments": severe_adjustments,
        "warnings": ([
            f"{len(severe_adjustments)} text boxes were shrunk below 75% of requested size"
        ] if severe_adjustments else []),
    }


def deck_slide_size(deck, page_entries):
    slide = deck.get("slide") or {}
    if not slide and page_entries:
        slide = page_entries[0]["manifest"].get("slide", {})
    return emu(slide.get("width", 13.333)), emu(slide.get("height", 7.5))


def write_deck(deck, page_entries, out_path, notes_entries):
    if not page_entries:
        raise ValueError("Deck has no pages")
    width, height = deck_slide_size(deck, page_entries)
    out = Path(out_path)
    out.parent.mkdir(parents=True, exist_ok=True)
    notes_by_page = {int(entry.get("page_index", 0)): entry for entry in notes_entries if entry.get("text")}
    notes_indices = sorted(notes_by_page)
    normalized_entries = [{**entry, "manifest": normalize_manifest(entry["manifest"])} for entry in page_entries]
    manifests = [entry["manifest"] for entry in normalized_entries]
    media_index = 1
    with zipfile.ZipFile(out, "w", zipfile.ZIP_DEFLATED) as z:
        z.writestr("[Content_Types].xml", content_types_xml(manifests, notes_indices))
        write_common_parts(z, len(page_entries), width, height, len(notes_by_page))
        for slide_index, entry in enumerate(normalized_entries, start=1):
            manifest = entry["manifest"]
            notes_index = slide_index if slide_index in notes_by_page else None
            z.writestr(f"ppt/slides/slide{slide_index}.xml", slide_xml(manifest))
            z.writestr(f"ppt/slides/_rels/slide{slide_index}.xml.rels", rels_xml(manifest, media_index, notes_index))
            base = Path(entry["manifest_path"]).resolve().parent
            for item in manifest.get("images", []):
                src = Path(item["path"])
                if not src.is_absolute():
                    src = base / src
                z.write(src, f"ppt/media/image{media_index}{image_ext(src)}")
                media_index += 1
            if notes_index is not None:
                note = notes_by_page[slide_index]
                notes_xml = note.get("notes_xml")
                if notes_xml and Path(notes_xml).exists():
                    z.writestr(f"ppt/notesSlides/notesSlide{notes_index}.xml", Path(notes_xml).read_bytes())
                else:
                    z.writestr(f"ppt/notesSlides/notesSlide{notes_index}.xml", notes_slide_xml(note.get("text", "")))
                z.writestr(f"ppt/notesSlides/_rels/notesSlide{notes_index}.xml.rels", notes_rels_xml(slide_index))
    normalize_powerpoint_package(out, len(page_entries), width, height, len(notes_by_page))


def page_entries_from_deck_manifest(deck_manifest_path):
    deck_path = Path(deck_manifest_path).resolve()
    deck = json.loads(deck_path.read_text(encoding="utf-8"))
    root = Path(deck.get("job_dir", deck_path.parent)).resolve()
    entries = []
    for page in deck.get("pages", []):
        manifest_path = Path(page.get("manifest", ""))
        if not manifest_path.is_absolute():
            manifest_path = root / manifest_path
        manifest = json.loads(manifest_path.read_text(encoding="utf-8"))
        entries.append({"manifest": manifest, "manifest_path": manifest_path})
    notes_path = deck.get("notes_manifest")
    notes_entries = []
    if notes_path:
        notes_file = Path(notes_path)
        if not notes_file.is_absolute():
            notes_file = root / notes_file
        if notes_file.exists():
            notes_entries = json.loads(notes_file.read_text(encoding="utf-8")).get("notes", [])
            for note in notes_entries:
                notes_xml = note.get("notes_xml")
                if notes_xml:
                    notes_xml_path = Path(notes_xml)
                    if not notes_xml_path.is_absolute():
                        notes_xml_path = root / notes_xml_path
                    note["notes_xml"] = str(notes_xml_path)
    return deck, entries, notes_entries


def output_path_from_deck_manifest(deck_manifest_path):
    deck_path = Path(deck_manifest_path).resolve()
    deck = json.loads(deck_path.read_text(encoding="utf-8"))
    root = Path(deck.get("job_dir", deck_path.parent)).resolve()
    output = Path(deck.get("output", "final/deck_edited.pptx"))
    if not output.is_absolute():
        output = root / output
    return output


def render_preview(manifest, manifest_path, out_path):
    from PIL import Image, ImageColor, ImageDraw, ImageFont

    manifest = normalize_manifest(manifest)
    slide = manifest.get("slide", {})
    width_in = float(slide.get("width", 13.333))
    height_in = float(slide.get("height", 7.5))
    scale = int(manifest.get("preview_scale", 120))
    canvas = Image.new("RGB", (int(width_in * scale), int(height_in * scale)), ImageColor.getrgb(slide.get("background", "#ffffff")))
    base = Path(manifest_path).resolve().parent
    draw = ImageDraw.Draw(canvas)

    def open_preview_image(src):
        if src.suffix.lower() != ".svg":
            return Image.open(src).convert("RGBA")
        commands = []
        for command in ("/opt/homebrew/bin/rsvg-convert", "/usr/local/bin/rsvg-convert"):
            if Path(command).exists():
                commands.append([command, "-o"])
        for command in ("/opt/homebrew/bin/magick", "/opt/homebrew/bin/convert"):
            if Path(command).exists():
                commands.append([command, str(src)])
        if not commands:
            print(f"Warning: cannot draft-preview SVG without rsvg-convert or ImageMagick: {src}", file=sys.stderr)
            return None
        with tempfile.NamedTemporaryFile(suffix=".png") as handle:
            errors = []
            for command in commands:
                invocation = [*command, handle.name, str(src)] if command[-1] == "-o" else [*command, handle.name]
                completed = subprocess.run(invocation, stdout=subprocess.PIPE, stderr=subprocess.PIPE)
                if completed.returncode == 0:
                    return Image.open(handle.name).convert("RGBA")
                errors.append(completed.stderr.decode("utf-8", errors="replace")[:200])
            print(f"Warning: SVG draft preview failed: {'; '.join(errors)}", file=sys.stderr)
            return None

    def render_shape(item):
        box = [item.get("left", 0) * scale, item.get("top", 0) * scale, (item.get("left", 0) + item.get("width", 1)) * scale, (item.get("top", 0) + item.get("height", 1)) * scale]
        fill = preview_color(item.get("fill"))
        outline = preview_color(item.get("stroke", "#000000"))
        width = max(1, int(float(item.get("stroke_width", 1))))
        if item.get("polygon"):
            points = [(point[0] * scale, point[1] * scale) for point in item["polygon"]]
            draw.polygon(points, fill=None if fill in (None, "none") else fill, outline=None if outline == "none" else outline)
        elif item.get("type") == "line":
            if "points" in item:
                points = [value * scale for value in item["points"]]
                draw.line(points, fill=outline, width=width)
                return
            if item.get("dash"):
                draw_dashed_line(draw, box, outline, width)
            else:
                draw.line(box, fill=outline, width=width)
        elif item.get("type") == "ellipse":
            draw.ellipse(box, fill=None if fill in (None, "none") else fill, outline=None if outline == "none" else outline, width=width)
        elif item.get("type") == "roundRect" or item.get("preset") == "roundRect":
            radius = int(float(item.get("radius", 0.12)) * scale)
            draw.rounded_rectangle(box, radius=radius, fill=None if fill in (None, "none") else fill, outline=None if outline == "none" else outline, width=width)
        elif item.get("preset") == "diamond":
            left, top, right, bottom = box
            center_x = (left + right) / 2
            center_y = (top + bottom) / 2
            points = [(center_x, top), (right, center_y), (center_x, bottom), (left, center_y)]
            draw.polygon(points, fill=None if fill in (None, "none") else fill, outline=None if outline == "none" else outline)
        else:
            draw.rectangle(box, fill=None if fill in (None, "none") else fill, outline=None if outline == "none" else outline, width=width)

    def render_image(item):
        src = Path(item["path"])
        if not src.is_absolute():
            src = base / src
        img = open_preview_image(src)
        if img is None:
            return
        img = img.resize((max(1, int(item.get("width", 1) * scale)), max(1, int(item.get("height", 1) * scale))))
        canvas.paste(img, (int(item.get("left", 0) * scale), int(item.get("top", 0) * scale)), img)

    def render_text(item):
        preview_font_scale = float(item.get("preview_font_scale", manifest.get("preview_font_scale", 1.0)))
        size = max(1, int(float(item.get("font_size", 18)) * scale / 72 * preview_font_scale))
        font_path = choose_preview_font(item.get("preview_font"))
        try:
            font = ImageFont.truetype(font_path, size=size) if font_path else ImageFont.load_default()
        except Exception:
            font = ImageFont.load_default()
        if item.get("paragraphs"):
            lines = []
            for paragraph in item["paragraphs"]:
                if isinstance(paragraph, str):
                    lines.append(paragraph)
                else:
                    lines.append("".join(str(run.get("text", "")) for run in paragraph.get("runs", [])))
            preview_text = "\n".join(lines)
        elif item.get("runs"):
            preview_text = "".join(str(run.get("text", "")) for run in item["runs"])
        else:
            preview_text = item.get("text", "")
        fill = preview_color(item.get("color", "#111111"))
        align = text_alignment(item.get("align", "left"))[1]
        valign = text_vertical_alignment(item.get("valign", "top"))[1]
        box_x = int(item.get("left", 0) * scale)
        box_y = int(item.get("top", 0) * scale)
        box_width = max(1, int(item.get("width", 1) * scale))
        box_height = max(1, int(item.get("height", 0.4) * scale))
        rotation = float(item.get("rotation", 0) or 0)

        def aligned_origin(bounds, origin_x, origin_y):
            text_width = bounds[2] - bounds[0]
            text_height = bounds[3] - bounds[1]
            if align == "center":
                text_x = origin_x + (box_width - text_width) // 2 - bounds[0]
            elif align == "right":
                text_x = origin_x + box_width - text_width - bounds[0]
            else:
                text_x = origin_x
            if valign == "middle":
                text_y = origin_y + (box_height - text_height) // 2 - bounds[1]
            elif valign == "bottom":
                text_y = origin_y + box_height - text_height - bounds[1]
            else:
                text_y = origin_y
            return text_x, text_y

        run_specs = []
        if item.get("runs"):
            cursor_x = 0
            base_size = size
            for run in item["runs"]:
                run_size = max(1, int(float(run.get("font_size", item.get("font_size", 18))) * scale / 72 * preview_font_scale))
                run_font_path = choose_preview_font(run.get("preview_font") or item.get("preview_font"))
                try:
                    run_font = ImageFont.truetype(run_font_path, size=run_size) if run_font_path else ImageFont.load_default()
                except Exception:
                    run_font = font
                run_fill = preview_color(run.get("color", item.get("color", "#111111")))
                baseline = float(run.get("baseline", 0) or 0)
                run_y = int(-baseline / 100000 * base_size)
                run_text = str(run.get("text", ""))
                run_specs.append((cursor_x, run_y, run_text, run_fill, run_font, run_font.getbbox(run_text)))
                cursor_x += int(round(draw.textlength(run_text, font=run_font)))

        def draw_content(target_draw, origin_x, origin_y):
            if run_specs:
                bounds = (
                    min(spec[0] + spec[5][0] for spec in run_specs),
                    min(spec[1] + spec[5][1] for spec in run_specs),
                    max(spec[0] + spec[5][2] for spec in run_specs),
                    max(spec[1] + spec[5][3] for spec in run_specs),
                )
                text_x, text_y = aligned_origin(bounds, origin_x, origin_y)
                for run_x, run_y, run_text, run_fill, run_font, _run_bounds in run_specs:
                    target_draw.text((text_x + run_x, text_y + run_y), run_text, fill=run_fill, font=run_font)
                return
            bounds = target_draw.multiline_textbbox((0, 0), preview_text, font=font, spacing=4, align=align)
            text_x, text_y = aligned_origin(bounds, origin_x, origin_y)
            target_draw.multiline_text((text_x, text_y), preview_text, fill=fill, font=font, spacing=4, align=align)

        if rotation:
            layer = Image.new("RGBA", (box_width, box_height), (0, 0, 0, 0))
            draw_content(ImageDraw.Draw(layer), 0, 0)
            rotated = layer.rotate(-rotation, expand=True)
            paste_x = box_x + (box_width - rotated.width) // 2
            paste_y = box_y + (box_height - rotated.height) // 2
            canvas.paste(rotated, (paste_x, paste_y), rotated)
            return
        draw_content(draw, box_x, box_y)

    def render_table(item):
        rows = item.get("rows") or []
        if not rows:
            return
        columns = max(len(row) for row in rows)
        left = int(item.get("left", 0) * scale)
        top = int(item.get("top", 0) * scale)
        width = max(1, int(item.get("width", 1) * scale))
        height = max(1, int(item.get("height", 1) * scale))
        column_weights = item.get("column_widths") or [1] * columns
        row_weights = item.get("row_heights") or [1] * len(rows)
        x_values = [left]
        for value in column_weights:
            x_values.append(x_values[-1] + int(width * float(value) / sum(column_weights)))
        x_values[-1] = left + width
        y_values = [top]
        for value in row_weights:
            y_values.append(y_values[-1] + int(height * float(value) / sum(row_weights)))
        y_values[-1] = top + height
        for row_index, row in enumerate(rows):
            for column_index in range(columns):
                value = row[column_index] if column_index < len(row) else ""
                cell = value if isinstance(value, dict) else {"text": value}
                fill = preview_color(cell.get("fill", item.get("header_fill" if row_index == 0 else "fill", "#FFFFFF")))
                border = preview_color(cell.get("border", item.get("border", "#B7C4D4")))
                box = [x_values[column_index], y_values[row_index], x_values[column_index + 1], y_values[row_index + 1]]
                draw.rectangle(box, fill=fill, outline=border)
                text = str(cell.get("text", ""))
                size = max(1, int(float(cell.get("font_size", item.get("font_size", 12))) * scale / 72))
                font_path = choose_preview_font(cell.get("preview_font") or item.get("preview_font"))
                try:
                    font = ImageFont.truetype(font_path, size=size) if font_path else ImageFont.load_default()
                except Exception:
                    font = ImageFont.load_default()
                bounds = draw.multiline_textbbox((0, 0), text, font=font, align="center")
                text_width = bounds[2] - bounds[0]
                text_height = bounds[3] - bounds[1]
                draw.multiline_text(
                    ((box[0] + box[2] - text_width) / 2, (box[1] + box[3] - text_height) / 2),
                    text,
                    font=font,
                    fill=preview_color(cell.get("color", item.get("header_color" if row_index == 0 else "color", "#111111"))),
                    align="center",
                )

    layered = []
    for index, item in enumerate(manifest.get("shapes", [])):
        layered.append((float(item.get("z_index", 100)), index, render_shape, item))
    for index, item in enumerate(manifest.get("images", [])):
        layered.append((float(item.get("z_index", 200)), index, render_image, item))
    for index, item in enumerate(manifest.get("tables", [])):
        layered.append((float(item.get("z_index", 250)), index, render_table, item))
    for index, item in enumerate(manifest.get("text_boxes", [])):
        layered.append((float(item.get("z_index", 300)), index, render_text, item))
    for _z_index, _order, renderer, item in sorted(layered, key=lambda entry: (entry[0], entry[1])):
        renderer(item)
    Path(out_path).parent.mkdir(parents=True, exist_ok=True)
    canvas.save(out_path)


def choose_preview_font(preferred):
    candidates = [
        preferred,
        "/System/Library/Fonts/STHeiti Medium.ttc",
        "/System/Library/Fonts/STHeiti Light.ttc",
        "/System/Library/Fonts/Supplemental/Arial Unicode.ttf",
        "/Library/Fonts/Arial Unicode.ttf",
    ]
    for candidate in candidates:
        if candidate and Path(candidate).exists():
            return candidate
    return None


def draw_dashed_line(draw, box, fill, width):
    x1, y1, x2, y2 = box
    dash = 8
    gap = 6
    if abs(y2 - y1) <= abs(x2 - x1):
        step = dash + gap
        x = min(x1, x2)
        end = max(x1, x2)
        y = y1
        while x < end:
            draw.line((x, y, min(x + dash, end), y), fill=fill, width=width)
            x += step
    else:
        step = dash + gap
        y = min(y1, y2)
        end = max(y1, y2)
        x = x1
        while y < end:
            draw.line((x, y, x, min(y + dash, end)), fill=fill, width=width)
            y += step


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("manifest", nargs="?")
    parser.add_argument("--deck-manifest")
    parser.add_argument("--out")
    parser.add_argument("--preview")
    parser.add_argument("--report")
    args = parser.parse_args()
    if args.deck_manifest:
        deck, entries, notes_entries = page_entries_from_deck_manifest(args.deck_manifest)
        out = Path(args.out) if args.out else output_path_from_deck_manifest(args.deck_manifest)
        write_deck(deck, entries, out, notes_entries)
        print(f"Wrote {out}")
        return
    if not args.manifest:
        parser.error("manifest is required unless --deck-manifest is used")
    if not args.out:
        parser.error("--out is required unless --deck-manifest provides an output")
    manifest = json.loads(Path(args.manifest).read_text(encoding="utf-8"))
    normalized = write_pptx(manifest, args.out, args.manifest)
    if args.report:
        report_path = Path(args.report)
        report_path.parent.mkdir(parents=True, exist_ok=True)
        report_path.write_text(
            json.dumps(build_report(normalized, args.out), ensure_ascii=False, indent=2) + "\n",
            encoding="utf-8",
        )
    if args.preview:
        render_preview(manifest, args.manifest, args.preview)
    print(f"Wrote {args.out}")
    if args.report:
        print(f"Wrote {args.report}")
    if args.preview:
        print(f"Wrote {args.preview}")


if __name__ == "__main__":
    main()
