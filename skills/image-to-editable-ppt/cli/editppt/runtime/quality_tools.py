#!/usr/bin/env python3
"""Deterministic PPTX readback, PowerPoint rendering, and image comparison.

These functions provide evidence to the page-owning Codex task. They do not
issue a product acceptance verdict.
"""

from __future__ import annotations

import hashlib
import json
import plistlib
import sys
from pathlib import Path
from typing import Any

import numpy as np
from PIL import Image, ImageChops
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from powerpoint_render import render_one_page
from editppt.text_evidence import region_text_coverage


POWERPOINT_APP = Path("/Applications/Microsoft PowerPoint.app")


def sha256(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as handle:
        for block in iter(lambda: handle.read(1024 * 1024), b""):
            digest.update(block)
    return digest.hexdigest()


def write_json(path: Path, value: dict[str, Any]) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(json.dumps(value, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")


def powerpoint_version() -> str:
    plist = POWERPOINT_APP / "Contents/Info.plist"
    try:
        value = plistlib.loads(plist.read_bytes())
    except (OSError, plistlib.InvalidFileException):
        return ""
    short = str(value.get("CFBundleShortVersionString") or "")
    build = str(value.get("CFBundleVersion") or "")
    return f"{short} ({build})" if short and build else short or build


def find_extractor(explicit: str = "") -> Path | None:
    """Return the Skill-owned target-only renderer for compatibility callers."""

    renderer = Path(__file__).with_name("powerpoint_render.py")
    return renderer.resolve() if renderer.is_file() else None


def render_powerpoint(
    pptx_path: Path,
    output_png: Path,
    *,
    evidence_dir: Path,
    extractor: str = "",
    dpi: int = 200,
) -> dict[str, Any]:
    """Render a deck through installed Microsoft PowerPoint, failing closed."""

    pptx_path = pptx_path.resolve()
    output_png = output_png.resolve()
    evidence_dir = evidence_dir.resolve()
    evidence_dir.mkdir(parents=True, exist_ok=True)
    if not pptx_path.is_file() or pptx_path.stat().st_size == 0:
        raise FileNotFoundError(f"non-empty PPTX not found: {pptx_path}")
    if sys.platform != "darwin" or not POWERPOINT_APP.is_dir():
        raise RuntimeError("Microsoft PowerPoint for macOS is unavailable")
    renderer = find_extractor(extractor)
    if renderer is None:
        raise RuntimeError("Skill-owned PowerPoint renderer is unavailable")
    candidate_sha = sha256(pptx_path)
    binding_path = evidence_dir / "render-binding.json"
    try:
        cached = json.loads(binding_path.read_text(encoding="utf-8"))
    except (OSError, json.JSONDecodeError):
        cached = {}
    if (
        cached.get("input_sha256") == candidate_sha
        and cached.get("renderer_version") == powerpoint_version()
        and int(cached.get("dpi") or 0) == int(dpi)
        and Path(str(cached.get("output_png") or "")).resolve() == output_png
        and output_png.is_file()
        and cached.get("output_sha256") == sha256(output_png)
    ):
        cached["cache_hit"] = True
        return cached
    report = render_one_page(
        pptx_path,
        output_png,
        evidence_dir=evidence_dir,
        dpi=dpi,
    )
    report_path = evidence_dir / "powerpoint-render.json"
    output_sha = sha256(output_png)
    bound = {
        "status": "rendered",
        "authoritative": True,
        "renderer": "microsoft-powerpoint",
        "renderer_version": powerpoint_version(),
        "input_pptx": str(pptx_path),
        "input_sha256": candidate_sha,
        "output_png": str(output_png),
        "output_sha256": output_sha,
        "dpi": dpi,
        "renderer_script": str(renderer),
        "renderer_report": str(report_path),
        "target_only": True,
        "canary_created": False,
        "cache_hit": False,
    }
    write_json(binding_path, bound)
    return bound


def _shape_text(shape: Any) -> str:
    if getattr(shape, "has_text_frame", False):
        return "\n".join(paragraph.text for paragraph in shape.text_frame.paragraphs)
    if getattr(shape, "has_table", False):
        return "\n".join(cell.text for row in shape.table.rows for cell in row.cells)
    return ""


def inspect_pptx(path: Path, text_hints_path: Path | None = None) -> dict[str, Any]:
    presentation = Presentation(str(path))
    slide_width = int(presentation.slide_width)
    slide_height = int(presentation.slide_height)
    pages: list[dict[str, Any]] = []
    candidate_regions: list[dict[str, Any]] = []
    total = {
        "object_count": 0,
        "text_shape_count": 0,
        "native_shape_count": 0,
        "table_count": 0,
        "picture_count": 0,
        "connector_count": 0,
        "out_of_bounds_count": 0,
        "max_picture_coverage": 0.0,
    }
    for page_index, slide in enumerate(presentation.slides, start=1):
        objects: list[dict[str, Any]] = []
        for shape in slide.shapes:
            left, top, width, height = map(int, (shape.left, shape.top, shape.width, shape.height))
            out = left < 0 or top < 0 or left + width > slide_width or top + height > slide_height
            picture = shape.shape_type == MSO_SHAPE_TYPE.PICTURE
            table = bool(getattr(shape, "has_table", False))
            connector = shape.shape_type == MSO_SHAPE_TYPE.LINE
            text = _shape_text(shape)
            coverage = (width * height) / max(1, slide_width * slide_height) if picture else 0.0
            record = {
                "name": shape.name,
                "type": str(shape.shape_type),
                "box_emu": [left, top, width, height],
                "text": text,
                "picture_coverage": round(coverage, 6),
                "out_of_bounds": out,
                "has_table": table,
            }
            objects.append(record)
            if text:
                candidate_regions.append({
                    "text": text,
                    "box_emu": [left, top, width, height],
                })
            total["object_count"] += 1
            total["text_shape_count"] += int(bool(text))
            total["picture_count"] += int(picture)
            total["table_count"] += int(table)
            total["connector_count"] += int(connector)
            total["native_shape_count"] += int(not picture)
            total["out_of_bounds_count"] += int(out)
            total["max_picture_coverage"] = max(total["max_picture_coverage"], coverage)
        pages.append({"page": page_index, "objects": objects})
    total["max_picture_coverage"] = round(float(total["max_picture_coverage"]), 6)
    payload = {
        "schema_version": 1,
        "pptx": str(path.resolve()),
        "sha256": sha256(path),
        "slide_count": len(presentation.slides),
        "slide_size_emu": [slide_width, slide_height],
        "summary": total,
        "risks": {
            "near_full_page_picture": total["max_picture_coverage"] >= 0.8,
            "out_of_bounds": total["out_of_bounds_count"] > 0,
        },
        "pages": pages,
    }
    if text_hints_path and text_hints_path.is_file() and len(presentation.slides) == 1:
        try:
            hints = json.loads(text_hints_path.read_text(encoding="utf-8"))
        except (OSError, json.JSONDecodeError):
            hints = {}
        source = hints.get("source") or {}
        source_width = float(source.get("width_px") or slide_width)
        source_height = float(source.get("height_px") or slide_height)
        normalized_candidates = []
        for region in candidate_regions:
            left, top, width, height = region["box_emu"]
            normalized_candidates.append({
                "text": region["text"],
                "box_px": [
                    left / slide_width * source_width,
                    top / slide_height * source_height,
                    width / slide_width * source_width,
                    height / slide_height * source_height,
                ],
            })
        expected = [
            {"text": str(line.get("text") or ""), "box_px": line.get("box_px")}
            for line in hints.get("lines", []) if isinstance(line, dict)
        ]
        evidence = region_text_coverage(expected, normalized_candidates)
        evidence["source"] = str(text_hints_path.resolve())
        evidence["advisory"] = True
        payload["text_evidence"] = evidence
        total["text_hint_coverage"] = evidence["text_coverage"]
        total["text_hint_missing_count"] = evidence["missing_text_count"]
    return payload


def compare_images(source: Path, candidate: Path, out_dir: Path) -> dict[str, Any]:
    out_dir.mkdir(parents=True, exist_ok=True)
    with Image.open(source).convert("RGB") as source_image, Image.open(candidate).convert("RGB") as candidate_image:
        candidate_fitted = candidate_image.resize(source_image.size, Image.Resampling.LANCZOS)
        source_array = np.asarray(source_image, dtype=np.float32) / 255.0
        candidate_array = np.asarray(candidate_fitted, dtype=np.float32) / 255.0
        delta = np.abs(source_array - candidate_array)
        rgb_loss = float(delta.mean())
        source_ink = np.mean(source_array, axis=2) < 0.93
        candidate_ink = np.mean(candidate_array, axis=2) < 0.93
        ink_loss = float(np.logical_xor(source_ink, candidate_ink).mean())
        heat = Image.fromarray(np.uint8(np.clip(delta.max(axis=2) * 4.0, 0, 1) * 255), "L")
        heat_rgb = Image.merge("RGB", (heat, Image.new("L", heat.size, 0), Image.new("L", heat.size, 0)))
        overlay = Image.blend(source_image, candidate_fitted, 0.5)
        diff = ImageChops.difference(source_image, candidate_fitted)
        heat_path = out_dir / "heatmap.png"
        overlay_path = out_dir / "overlay.png"
        diff_path = out_dir / "diff.png"
        heat_rgb.save(heat_path)
        overlay.save(overlay_path)
        diff.save(diff_path)
        # Cluster meaningful differences on a compact tile grid.  These
        # source/candidate pairs let Codex inspect only the regions that
        # changed instead of re-reading the entire page after every repair.
        difference_mask = delta.max(axis=2) > 0.12
        rows, columns = 18, 32
        tile_h = max(1, int(np.ceil(source_image.height / rows)))
        tile_w = max(1, int(np.ceil(source_image.width / columns)))
        active: set[tuple[int, int]] = set()
        for row in range(rows):
            for column in range(columns):
                tile = difference_mask[
                    row * tile_h : min(source_image.height, (row + 1) * tile_h),
                    column * tile_w : min(source_image.width, (column + 1) * tile_w),
                ]
                if tile.size and float(tile.mean()) >= 0.04:
                    active.add((row, column))
        components: list[list[tuple[int, int]]] = []
        while active:
            stack = [active.pop()]
            component = []
            while stack:
                current = stack.pop()
                component.append(current)
                row, column = current
                for neighbor in ((row - 1, column), (row + 1, column), (row, column - 1), (row, column + 1)):
                    if neighbor in active:
                        active.remove(neighbor)
                        stack.append(neighbor)
            components.append(component)
        components.sort(key=len, reverse=True)
        region_dir = out_dir / "regions"
        region_dir.mkdir(parents=True, exist_ok=True)
        for stale in region_dir.glob("*.png"):
            stale.unlink()
        regions = []
        for index, component in enumerate(components[:12], start=1):
            region_rows = [value[0] for value in component]
            region_columns = [value[1] for value in component]
            left = max(0, min(region_columns) * tile_w - tile_w // 2)
            top = max(0, min(region_rows) * tile_h - tile_h // 2)
            right = min(source_image.width, (max(region_columns) + 1) * tile_w + tile_w // 2)
            bottom = min(source_image.height, (max(region_rows) + 1) * tile_h + tile_h // 2)
            box = (left, top, right, bottom)
            source_crop = region_dir / f"region-{index:02d}-source.png"
            candidate_crop = region_dir / f"region-{index:02d}-candidate.png"
            diff_crop = region_dir / f"region-{index:02d}-diff.png"
            source_image.crop(box).save(source_crop)
            candidate_fitted.crop(box).save(candidate_crop)
            diff.crop(box).save(diff_crop)
            regions.append({
                "id": f"region-{index:02d}",
                "box_px": [left, top, right - left, bottom - top],
                "tile_count": len(component),
                "source": str(source_crop),
                "candidate": str(candidate_crop),
                "diff": str(diff_crop),
            })
    report = {
        "schema_version": 1,
        "source": str(source.resolve()),
        "candidate": str(candidate.resolve()),
        "source_sha256": sha256(source),
        "candidate_sha256": sha256(candidate),
        "coarse_rgb_loss": round(rgb_loss, 6),
        "content_ink_loss": round(ink_loss, 6),
        "diff": str(diff_path),
        "heatmap": str(heat_path),
        "overlay": str(overlay_path),
        "regions": regions,
        "note": "Diagnostic evidence only; inspect region geometry and text directly.",
    }
    write_json(out_dir / "compare.json", report)
    return report
