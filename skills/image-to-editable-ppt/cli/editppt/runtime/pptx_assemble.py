#!/usr/bin/env python3
"""Relationship-aware assembly of independent one-slide PPTX packages."""

from __future__ import annotations

import json
import posixpath
import shutil
import tempfile
import zipfile
from copy import deepcopy
from pathlib import Path
from typing import Any
from xml.etree import ElementTree as ET

from pptx import Presentation

from quality_tools import sha256, write_json


REL_NS = "http://schemas.openxmlformats.org/package/2006/relationships"
P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
R_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
CT_NS = "http://schemas.openxmlformats.org/package/2006/content-types"
REL_TAG = f"{{{REL_NS}}}Relationship"


def _rels_name(part: str) -> str:
    directory, name = posixpath.split(part)
    return posixpath.join(directory, "_rels", name + ".rels")


def _resolve(source_part: str, target: str) -> str:
    return posixpath.normpath(posixpath.join(posixpath.dirname(source_part), target))


def _relative(source_part: str, target_part: str) -> str:
    return posixpath.relpath(target_part, posixpath.dirname(source_part))


def _source_slide(package: dict[str, bytes]) -> str:
    presentation = ET.fromstring(package["ppt/presentation.xml"])
    rels = ET.fromstring(package["ppt/_rels/presentation.xml.rels"])
    relation_by_id = {value.get("Id"): value for value in rels.findall(REL_TAG)}
    slide_ids = presentation.findall(f".//{{{P_NS}}}sldId")
    if len(slide_ids) != 1:
        raise ValueError(f"page PPTX must contain one slide; found {len(slide_ids)}")
    rel_id = slide_ids[0].get(f"{{{R_NS}}}id")
    relation = relation_by_id.get(rel_id)
    if relation is None:
        raise ValueError("page PPTX slide relationship is missing")
    return _resolve("ppt/presentation.xml", str(relation.get("Target") or ""))


def _content_types(package: dict[str, bytes]) -> tuple[dict[str, str], dict[str, str]]:
    root = ET.fromstring(package["[Content_Types].xml"])
    defaults: dict[str, str] = {}
    overrides: dict[str, str] = {}
    for child in root:
        if child.tag == f"{{{CT_NS}}}Default":
            defaults[str(child.get("Extension") or "").lower()] = str(child.get("ContentType") or "")
        elif child.tag == f"{{{CT_NS}}}Override":
            overrides[str(child.get("PartName") or "").lstrip("/")] = str(child.get("ContentType") or "")
    return defaults, overrides


def _part_content_type(part: str, defaults: dict[str, str], overrides: dict[str, str]) -> str:
    if part in overrides:
        return overrides[part]
    extension = posixpath.splitext(part)[1].lstrip(".").lower()
    return defaults.get(extension, "")


class PackageMerger:
    def __init__(self, base: dict[str, bytes]) -> None:
        self.base = base
        self.base_defaults, self.base_overrides = _content_types(base)
        self.added_defaults: dict[str, str] = {}
        self.added_overrides: dict[str, str] = {}
        self.counters: dict[tuple[str, str, str], int] = {}
        self.base_layout = self._base_layout_part()

    def _base_layout_part(self) -> str:
        slide = _source_slide(self.base)
        rels_name = _rels_name(slide)
        rels = ET.fromstring(self.base[rels_name])
        for relation in rels.findall(REL_TAG):
            if str(relation.get("Type") or "").endswith("/slideLayout"):
                return _resolve(slide, str(relation.get("Target") or ""))
        raise ValueError("base page PPTX has no slide layout relationship")

    def unique_part(self, requested: str) -> str:
        if requested not in self.base:
            return requested
        directory, name = posixpath.split(requested)
        stem, extension = posixpath.splitext(name)
        key = (directory, stem, extension)
        value = self.counters.get(key, 2)
        while True:
            candidate = posixpath.join(directory, f"{stem}_editppt_{value}{extension}")
            if candidate not in self.base:
                self.counters[key] = value + 1
                return candidate
            value += 1

    def register_content_type(
        self,
        source_part: str,
        target_part: str,
        defaults: dict[str, str],
        overrides: dict[str, str],
    ) -> None:
        content_type = _part_content_type(source_part, defaults, overrides)
        if not content_type:
            return
        extension = posixpath.splitext(target_part)[1].lstrip(".").lower()
        if source_part in overrides:
            self.added_overrides[target_part] = content_type
        elif extension and extension not in self.base_defaults:
            self.added_defaults[extension] = content_type

    def copy_graph(
        self,
        source_package: dict[str, bytes],
        source_part: str,
        target_part: str,
        *,
        source_slide: str,
        target_slide: str,
        mapping: dict[str, str],
        defaults: dict[str, str],
        overrides: dict[str, str],
    ) -> str:
        if source_part in mapping:
            return mapping[source_part]
        mapping[source_part] = target_part
        if source_part not in source_package:
            raise ValueError(f"PPTX relationship target is missing: {source_part}")
        self.base[target_part] = source_package[source_part]
        self.register_content_type(source_part, target_part, defaults, overrides)
        source_rels_name = _rels_name(source_part)
        if source_rels_name not in source_package:
            return target_part
        rels = ET.fromstring(source_package[source_rels_name])
        for relation in rels.findall(REL_TAG):
            if relation.get("TargetMode") == "External":
                continue
            relation_target = str(relation.get("Target") or "")
            source_target = _resolve(source_part, relation_target)
            relation_type = str(relation.get("Type") or "")
            if target_part == target_slide and relation_type.endswith("/slideLayout"):
                # Page candidates must carry their visible background and
                # objects on the slide. Reusing the base blank layout avoids
                # importing a second master/theme graph, which PowerPoint can
                # open but may refuse to export after a package-level merge.
                copied_target = self.base_layout
            elif relation_type.endswith("/slide") and source_target == source_slide:
                copied_target = target_slide
            elif source_target in mapping:
                copied_target = mapping[source_target]
            else:
                requested = source_target
                copied_target = self.unique_part(requested)
                copied_target = self.copy_graph(
                    source_package,
                    source_target,
                    copied_target,
                    source_slide=source_slide,
                    target_slide=target_slide,
                    mapping=mapping,
                    defaults=defaults,
                    overrides=overrides,
                )
            relation.set("Target", _relative(target_part, copied_target))
        target_rels_name = _rels_name(target_part)
        self.base[target_rels_name] = ET.tostring(rels, encoding="UTF-8", xml_declaration=True)
        return target_part

    def finalize_types(self) -> None:
        root = ET.fromstring(self.base["[Content_Types].xml"])
        for extension, content_type in sorted(self.added_defaults.items()):
            if extension not in self.base_defaults:
                ET.SubElement(root, f"{{{CT_NS}}}Default", Extension=extension, ContentType=content_type)
                self.base_defaults[extension] = content_type
        for part, content_type in sorted(self.added_overrides.items()):
            if part not in self.base_overrides:
                ET.SubElement(root, f"{{{CT_NS}}}Override", PartName="/" + part, ContentType=content_type)
                self.base_overrides[part] = content_type
        ET.register_namespace("", CT_NS)
        self.base["[Content_Types].xml"] = ET.tostring(root, encoding="UTF-8", xml_declaration=True)


def _read_package(path: Path) -> dict[str, bytes]:
    with zipfile.ZipFile(path) as package:
        bad = package.testzip()
        if bad:
            raise ValueError(f"corrupt PPTX member: {bad}")
        return {name: package.read(name) for name in package.namelist()}


def _slide_size(package: dict[str, bytes]) -> tuple[int, int]:
    presentation = ET.fromstring(package["ppt/presentation.xml"])
    value = presentation.find(f"{{{P_NS}}}sldSz")
    if value is None:
        raise ValueError("PPTX slide size is missing")
    return int(value.get("cx") or 0), int(value.get("cy") or 0)


def _append_slide_reference(base: dict[str, bytes], slide_part: str) -> None:
    presentation = ET.fromstring(base["ppt/presentation.xml"])
    relationships = ET.fromstring(base["ppt/_rels/presentation.xml.rels"])
    slide_list = presentation.find(f"{{{P_NS}}}sldIdLst")
    if slide_list is None:
        slide_list = ET.SubElement(presentation, f"{{{P_NS}}}sldIdLst")
    existing_ids = [int(value.get("id") or 255) for value in slide_list]
    existing_rel_ids = {str(value.get("Id") or "") for value in relationships.findall(REL_TAG)}
    rel_number = 1
    while f"rId{rel_number}" in existing_rel_ids:
        rel_number += 1
    rel_id = f"rId{rel_number}"
    ET.SubElement(
        relationships,
        REL_TAG,
        Id=rel_id,
        Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide",
        Target=_relative("ppt/presentation.xml", slide_part),
    )
    slide_id = ET.SubElement(slide_list, f"{{{P_NS}}}sldId")
    slide_id.set("id", str(max(existing_ids, default=255) + 1))
    slide_id.set(f"{{{R_NS}}}id", rel_id)
    ET.register_namespace("p", P_NS)
    ET.register_namespace("r", R_NS)
    base["ppt/presentation.xml"] = ET.tostring(presentation, encoding="UTF-8", xml_declaration=True)
    base["ppt/_rels/presentation.xml.rels"] = ET.tostring(relationships, encoding="UTF-8", xml_declaration=True)


def _update_app_slide_count(base: dict[str, bytes], count: int) -> None:
    name = "docProps/app.xml"
    if name not in base:
        return
    root = ET.fromstring(base[name])
    slides = next((value for value in root if value.tag.endswith("}Slides")), None)
    if slides is not None:
        slides.text = str(count)
        base[name] = ET.tostring(root, encoding="UTF-8", xml_declaration=True)


def assemble_pptx_packages(inputs: list[Path], out: Path, evidence_dir: Path) -> dict[str, Any]:
    if not inputs:
        raise ValueError("at least one page PPTX is required")
    inputs = [value.resolve() for value in inputs]
    for value in inputs:
        if not value.is_file() or value.stat().st_size == 0:
            raise FileNotFoundError(f"non-empty page PPTX not found: {value}")
        if len(Presentation(str(value)).slides) != 1:
            raise ValueError(f"page input must contain exactly one slide: {value}")
    out = out.resolve()
    evidence_dir.mkdir(parents=True, exist_ok=True)
    if len(inputs) == 1:
        out.parent.mkdir(parents=True, exist_ok=True)
        shutil.copy2(inputs[0], out)
    else:
        base = _read_package(inputs[0])
        base_size = _slide_size(base)
        merger = PackageMerger(base)
        next_slide = 2
        for source_path in inputs[1:]:
            source = _read_package(source_path)
            if _slide_size(source) != base_size:
                raise ValueError(f"slide size mismatch: {source_path}")
            source_slide = _source_slide(source)
            while f"ppt/slides/slide{next_slide}.xml" in base:
                next_slide += 1
            target_slide = f"ppt/slides/slide{next_slide}.xml"
            defaults, overrides = _content_types(source)
            merger.copy_graph(
                source,
                source_slide,
                target_slide,
                source_slide=source_slide,
                target_slide=target_slide,
                mapping={},
                defaults=defaults,
                overrides=overrides,
            )
            _append_slide_reference(base, target_slide)
            next_slide += 1
        merger.finalize_types()
        _update_app_slide_count(base, len(inputs))
        out.parent.mkdir(parents=True, exist_ok=True)
        with tempfile.NamedTemporaryFile(dir=out.parent, suffix=".pptx", delete=False) as handle:
            temporary = Path(handle.name)
        try:
            with zipfile.ZipFile(temporary, "w", zipfile.ZIP_DEFLATED) as package:
                for name, data in base.items():
                    package.writestr(name, data)
            assembled = Presentation(str(temporary))
            if len(assembled.slides) != len(inputs):
                raise RuntimeError(f"assembled {len(assembled.slides)} slides, expected {len(inputs)}")
            temporary.replace(out)
        finally:
            temporary.unlink(missing_ok=True)
    result = {
        "status": "ready",
        "assembler": "relationship-aware-ooxml",
        "inputs": [str(value) for value in inputs],
        "input_sha256": [sha256(value) for value in inputs],
        "output": str(out),
        "output_sha256": sha256(out),
        "slide_count": len(Presentation(str(out)).slides),
    }
    write_json(evidence_dir / "assemble.json", result)
    return result
